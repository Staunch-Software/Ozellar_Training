"""FastAPI application for the Ozellar Marine seafarer training portal.

Auth: crew log in with crew_id + date of birth (DDMMYYYY); admins with
email + password. All learner data (progress, assessments, certificates)
is tied to the authenticated user via a JWT bearer token.

Run:  uvicorn app.main:app --reload
"""
import csv
import io
import os
import re
import shutil
import subprocess
import tempfile
import uuid
from datetime import datetime, timezone, timedelta
from typing import Optional

from fastapi import FastAPI, Depends, HTTPException, UploadFile, File, Form, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import Response, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from sqlalchemy.orm import Session
from sqlalchemy import func

from .database import get_db, SessionLocal
from . import models, email_service, storage
from .video import compress_video
from .certificates import build_certificate_pdf
from .auth import (
    create_token, get_current_user, require_admin, user_public,
    verify_password, hash_password, parse_ddmmyyyy, normalize_name,
    check_rate_limit, clear_rate_limit,
    # screening test auth
    create_screening_token, get_current_candidate, candidate_public,
    bearer, SECRET_KEY, ALGORITHM,
)

# public origin used in the certificate's verification line / verify links
PUBLIC_BASE_URL = os.getenv("PUBLIC_BASE_URL", "https://training.ozellar.com")

# admin course-builder uploads (slide images rendered from .pptx, uploaded
# videos) — served back out under /api/uploads so nginx's existing /api
# proxy (see deploy/nginx-ozellar.conf) covers it with no further config.
UPLOAD_DIR = os.getenv(
    "UPLOAD_DIR", os.path.join(os.path.dirname(__file__), "..", "uploads"))
os.makedirs(UPLOAD_DIR, exist_ok=True)

from contextlib import asynccontextmanager
from apscheduler.schedulers.asyncio import AsyncIOScheduler
import jwt

smartpal_scheduler = None
email_scheduler = None

def send_pending_digest_job():
    # Job to scan for pending approvals and send the digest
    from datetime import datetime, timezone, timedelta
    
    admin_email = os.getenv("ADMIN_EMAIL")
    if not admin_email:
        print("[send_pending_digest] Skipped: ADMIN_EMAIL not set")
        return

    with SessionLocal() as db:
        pending = db.query(models.AssessmentApproval).filter_by(status="pending", digest_sent=False).all()
        if not pending:
            return
            
        approvals_list = []
        for ap in pending:
            user = db.get(models.User, ap.learner_id)
            course = db.get(models.Course, ap.course_id)
            if not user or not course:
                continue
                
            # Create a one-off token for approval. Valid for 7 days.
            from .auth import SECRET_KEY, ALGORITHM
            token_payload = {
                "sub": f"approve:{ap.id}",
                "type": "approval",
                "exp": datetime.now(timezone.utc) + timedelta(days=7)
            }
            token = jwt.encode(token_payload, SECRET_KEY, algorithm=ALGORITHM)
            ap.approval_token = token
            ap.digest_sent = True
            
            approvals_list.append({
                "learner_name": user.full_name,
                "crew_id": user.crew_id,
                "course_title": course.title,
                "score": ap.score,
                "token": token
            })
            
        db.commit()
        if approvals_list:
            email_service.send_digest_email(admin_email, approvals_list)


@asynccontextmanager
async def lifespan(app: FastAPI):
    global smartpal_scheduler, email_scheduler

    # 1. DB auto-seed
    db = SessionLocal()
    try:
        if db.query(models.Course).count() == 0 or db.query(models.User).count() == 0:
            from .seed import run
            run()
    except Exception as e:
        print(f"[startup] auto-seed skipped ({e}). "
              f"Run:  alembic upgrade head  &&  python -m app.seed")
    finally:
        db.close()

    # 2. Start Schedulers
    from apscheduler.schedulers.asyncio import AsyncIOScheduler
    from .smartpal_sync import schedule_jobs
    import asyncio
    
    if os.getenv("SMARTPAL_USERNAME") and os.getenv("SMARTPAL_PASSWORD"):
        smartpal_scheduler = AsyncIOScheduler()
        schedule_jobs(smartpal_scheduler)
        smartpal_scheduler.start()
        print("[startup] SmartPAL sync scheduled for 7:00 and 19:00 IST")
    else:
        print("[startup] SmartPAL sync disabled (SMARTPAL_USERNAME/PASSWORD not set)")

    interval = int(os.getenv("APPROVAL_DIGEST_INTERVAL_MINUTES", "30"))
    email_scheduler = AsyncIOScheduler()
    email_scheduler.add_job(send_pending_digest_job, 'interval', minutes=interval)
    email_scheduler.start()
    print(f"[startup] Approval digest email scheduled every {interval} minutes")

    try:
        yield
    except asyncio.CancelledError:
        pass
    finally:
        # Shutdown Schedulers
        if smartpal_scheduler:
            smartpal_scheduler.shutdown(wait=False)
        if email_scheduler:
            email_scheduler.shutdown(wait=False)

app = FastAPI(title="Ozellar Marine Training API", version="0.2.0", lifespan=lifespan)

app.add_middleware(
    CORSMiddleware,
    allow_origins=os.getenv("CORS_ORIGINS", "http://localhost:5173").split(","),
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# NOTE: We do NOT use StaticFiles for uploads because it does not support
# HTTP Range requests, which are required for HTML5 video seeking.
# Instead, we use a custom endpoint below that handles Range headers properly.

import mimetypes

@app.get("/api/uploads/{course_id}/{filename}")
async def serve_upload(course_id: str, filename: str, request: Request):
    file_size = storage.get_size(course_id, filename)
    if file_size is None:
        raise HTTPException(status_code=404, detail="File not found")

    mime_type, _ = mimetypes.guess_type(filename)
    if not mime_type:
        mime_type = "application/octet-stream"

    range_header = request.headers.get("range")
    if range_header:
        # Parse Range: bytes=start-end
        range_match = re.match(r"bytes=(\d+)-(\d*)", range_header)
        if range_match:
            start = int(range_match.group(1))
            end = int(range_match.group(2)) if range_match.group(2) else file_size - 1
            end = min(end, file_size - 1)
            chunk_size = end - start + 1

            return StreamingResponse(
                storage.stream_range(course_id, filename, start, end),
                status_code=206,
                media_type=mime_type,
                headers={
                    "Content-Range": f"bytes {start}-{end}/{file_size}",
                    "Accept-Ranges": "bytes",
                    "Content-Length": str(chunk_size),
                    "Cache-Control": "public, max-age=3600",
                }
            )

    # No range header - serve entire file
    return StreamingResponse(
        storage.stream_full(course_id, filename),
        media_type=mime_type,
        headers={
            "Accept-Ranges": "bytes",
            "Content-Length": str(file_size),
            "Cache-Control": "public, max-age=3600",
        }
    )


# ======================= AUTH =======================
class LoginRequest(BaseModel):
    mode: str                       # 'crew' | 'admin'
    name: str | None = None         # crew login: full name
    dob: str | None = None          # 8 digits DDMMYYYY
    crewId: str | None = None       # crew login: tiebreaker only (name+DOB collision)
    email: str | None = None
    password: str | None = None


@app.post("/api/auth/login")
def login(req: LoginRequest, db: Session = Depends(get_db)):
    if req.mode == "crew":
        name = (req.name or "").strip()
        dob_raw = (req.dob or "").strip()
        check_rate_limit(db, f"crew:{normalize_name(name)}:{dob_raw}")
        dob = parse_ddmmyyyy(dob_raw)
        if not name or dob is None:
            raise HTTPException(401, "Invalid name or date of birth")
        # Match on date of birth (narrows to a few rows), then compare the
        # normalized name in Python so case/spacing differences don't matter.
        candidates = db.query(models.User).filter_by(role="learner", date_of_birth=dob).all()
        target = normalize_name(name)
        matches = [u for u in candidates if normalize_name(u.full_name) == target]
        if not matches:
            raise HTTPException(401, "Invalid name or date of birth")
        if len(matches) > 1:
            # rare collision: two crew share name + DOB → need Crew ID to disambiguate
            crew_id = (req.crewId or "").strip()
            if not crew_id:
                raise HTTPException(
                    409, "More than one crew member matches that name and date of birth. "
                         "Please enter your Crew ID to continue.")
            matches = [u for u in matches if u.crew_id == crew_id]
            if len(matches) != 1:
                raise HTTPException(401, "Invalid Crew ID for that name and date of birth")
        user = matches[0]
        clear_rate_limit(db, f"crew:{normalize_name(name)}:{dob_raw}")

    elif req.mode == "admin":
        email = (req.email or "").strip().lower()
        check_rate_limit(db, f"admin:{email}")
        user = db.query(models.User).filter_by(email=email, role="admin").first()
        if not user or not user.password_hash or not verify_password(req.password or "", user.password_hash):
            raise HTTPException(401, "Invalid email or password")
        clear_rate_limit(db, f"admin:{email}")
    else:
        raise HTTPException(400, "Invalid login mode")

    if not user.is_active:
        raise HTTPException(403, "This account is disabled")
    return {"token": create_token(user), "user": user_public(user)}


@app.get("/api/auth/crew-search")
def crew_search(q: str, request: Request, db: Session = Depends(get_db)):
    """Public (pre-login) name autocomplete for the crew sign-in form.
    Deliberately minimal: only name + rank (never crew_id/DOB/passport/etc.),
    active learners only, capped result count, and a per-IP rate limit
    separate from the stricter login-attempt limiter — this endpoint has no
    credential to check, just a lookup, so it needs its own looser budget
    that still blocks bulk roster scraping."""
    check_rate_limit(db, f"crew-search:{request.client.host if request.client else 'unknown'}",
                     max_attempts=40, window_seconds=60)
    query = normalize_name(q)
    if len(query) < 1:
        return []
    candidates = (db.query(models.User)
                  .filter_by(role="learner", is_active=True).all())
    matches = [u for u in candidates if query in normalize_name(u.full_name)]
    matches.sort(key=lambda u: (not normalize_name(u.full_name).startswith(query), u.full_name))
    return [{"name": u.full_name, "rank": u.rank} for u in matches[:8]]


@app.get("/api/auth/me")
def me(
    creds = Depends(bearer),
    db: Session = Depends(get_db),
):
    """Universal /me endpoint — handles session tokens (crew/admin) and
    test_session tokens (screening test candidates)."""
    if creds is None:
        raise HTTPException(401, "Not authenticated")
    try:
        payload = jwt.decode(creds.credentials, SECRET_KEY, algorithms=[ALGORITHM])
    except jwt.ExpiredSignatureError:
        raise HTTPException(401, "Session expired — please sign in again")
    except jwt.PyJWTError:
        raise HTTPException(401, "Invalid authentication token")

    token_type = payload.get("type", "session")
    if token_type == "session":
        user = db.get(models.User, str(payload["sub"]))
        if not user or not user.is_active:
            raise HTTPException(401, "User not found or inactive")
        return user_public(user)
    elif token_type == "test_session":
        candidate = db.get(models.ScreeningCandidate, str(payload["sub"]))
        if not candidate or not candidate.is_active:
            raise HTTPException(401, "Candidate not found or inactive")
        return candidate_public(candidate)
    else:
        raise HTTPException(401, "Invalid token type")



# ======================= COURSES =======================
def get_progress(db, learner_id, course_id):
    return (db.query(models.Progress)
            .filter_by(learner_id=learner_id, course_id=course_id).first())


def notify(db, user_id, kind, title, body=None, link=None):
    """Create an in-app notification (caller commits)."""
    db.add(models.Notification(user_id=user_id, kind=kind, title=title,
                               body=body, link=link))


def enrolled_course_ids(db, user_id):
    return [e.course_id for e in
            db.query(models.Enrollment).filter_by(learner_id=user_id).all()]


def require_enrollment(db, user, course_id):
    """Admins may access any course (preview); learners must be assigned it."""
    if user.role == "admin":
        return
    is_enrolled = db.query(models.Enrollment).filter_by(
        learner_id=user.id, course_id=course_id).first()
    if not is_enrolled:
        raise HTTPException(403, "You are not assigned to this course")


def serialize_course(db, course, progress, detail=False):
    done = set(progress.completed_chapters or []) if progress else set()
    total = len(course.chapters)
    completed = sum(1 for ch in course.chapters if ch.id in done)
    pct = round(completed / total * 100) if total else 0
    data = {
        "id": course.id, "slug": course.slug, "title": course.title,
        "subtitle": course.subtitle, "icon": course.icon, "gradient": course.gradient,
        "durationLabel": course.duration_label, "status": course.status,
        "statusNote": course.status_note, "total": total, "completedCount": completed,
        "progressPct": pct,
        "passed": bool(progress.passed) if progress else False,
        "score": progress.score if progress else None,
        "hasAssessment": len(course.questions) > 0,
    }
    
    # Check if certificate is pending
    cert_pending = False
    if progress and progress.passed:
        cert_pending = db.query(models.AssessmentApproval).filter_by(
            learner_id=progress.learner_id, course_id=course.id, status="pending"
        ).first() is not None
    data["certPending"] = cert_pending
    if detail:
        attempts_used = 0
        if progress:
            attempts_used = db.query(models.Attempt).filter_by(
                learner_id=progress.learner_id, course_id=course.id
            ).count()
        data["chapters"] = [{
            "id": ch.id, "n": ch.n, "chapterLabel": ch.chapter_label, "title": ch.title,
            "intro": ch.intro, "sections": ch.sections, "figure": ch.figure,
            "image": ch.image, "videos": ch.videos, "done": ch.id in done,
            "kind": ch.kind,
            # checkpoint quizzes are ungraded (non-blocking) — safe to send
            # the answer key straight to the client, unlike the final assessment
            "quizQuestions": [{
                "q": qq.prompt, "options": qq.options, "answer": qq.answer, "explain": qq.explain,
            } for qq in ch.quiz_questions] if ch.kind == "quiz" else [],
        } for ch in course.chapters]
        data["cert"] = course.cert
        # Integrity: the answer key and explanations are NOT sent to the client.
        # Grading is server-side (submit_assessment) and the correct answers +
        # explanations come back only in the graded result.
        data["assessment"] = {
            "passMark": course.pass_mark,
            "maxAttempts": course.max_attempts,
            "attemptsUsed": attempts_used,
            "questions": [{
                "q": q.prompt, "options": q.options,
            } for q in course.questions],
        }
        
        # Include latest attempt if exists so the frontend can resume/show results directly
        if progress:
            latest = db.query(models.Attempt).filter_by(
                learner_id=progress.learner_id, course_id=course.id
            ).order_by(models.Attempt.created_at.desc()).first()
            if latest:
                can_retry = False
                if not latest.passed:
                    if course.max_attempts is None:
                        can_retry = True
                    else:
                        can_retry = attempts_used < course.max_attempts
                
                show_answers = latest.passed or not can_retry
                questions = sorted(course.questions, key=lambda q: q.order)
                review = [{
                    "q": q.prompt, "options": q.options,
                    "correct": q.answer if (show_answers or a == q.answer) else None,
                    "chosen": a, "isCorrect": a == q.answer,
                    "explain": q.explain if (show_answers or a == q.answer) else None,
                } for q, a in zip(questions, latest.answers)]
                
                data["latestAttempt"] = {
                    "score": latest.score,
                    "passed": latest.passed,
                    "correct": sum(1 for q, a in zip(questions, latest.answers) if a == q.answer),
                    "total": len(questions),
                    "review": review,
                    "attemptsUsed": attempts_used,
                    "maxAttempts": course.max_attempts,
                    "certPending": cert_pending
                }
    else:
        data["chapters"] = [{"id": ch.id} for ch in course.chapters]
    return data


@app.get("/api/learner")
def learner(user: models.User = Depends(get_current_user)):
    return user_public(user)


@app.get("/api/courses")
def list_courses(user: models.User = Depends(get_current_user), db: Session = Depends(get_db)):
    lid = user.id
    q = db.query(models.Course).order_by(models.Course.order)
    if user.role != "admin":
        # learners see only the courses assigned to them
        assigned = enrolled_course_ids(db, user.id)
        if not assigned:
            return []
        q = q.filter(models.Course.id.in_(assigned))
    return [serialize_course(db, c, get_progress(db, lid, c.id)) for c in q.all()]


@app.get("/api/courses/{slug}")
def get_course(slug: str, user: models.User = Depends(get_current_user), db: Session = Depends(get_db)):
    c = db.query(models.Course).filter_by(slug=slug).first()
    if not c:
        raise HTTPException(404, "Course not found")
    require_enrollment(db, user, c.id)
    return serialize_course(db, c, get_progress(db, user.id, c.id), detail=True)


@app.post("/api/courses/{course_id}/chapters/{chapter_id}/complete")
def complete_chapter(course_id: str, chapter_id: str,
                     user: models.User = Depends(get_current_user), db: Session = Depends(get_db)):
    require_enrollment(db, user, course_id)
    lid = user.id
    course = db.query(models.Course).filter_by(id=course_id).first()
    if not course:
        raise HTTPException(404, "Course not found")
    p = get_progress(db, lid, course_id)
    if not p:
        p = models.Progress(learner_id=lid, course_id=course_id, completed_chapters=[])
        db.add(p)
    done = list(p.completed_chapters or [])
    if chapter_id not in done:
        done.append(chapter_id)
    p.completed_chapters = done

    # Auto-complete courses that have no final assessment once all chapters are done
    auto_completed = False
    if not course.questions:  # no assessment questions
        all_chapter_ids = {ch.id for ch in course.chapters}
        if all_chapter_ids and set(done) >= all_chapter_ids and not p.passed:
            p.passed = True
            db.flush()  # ensure p.id is available

            # Queue for admin approval (same flow as a passed assessment)
            from .auth import SECRET_KEY, ALGORITHM
            from datetime import timedelta
            ap = models.AssessmentApproval(
                learner_id=user.id, course_id=course.id,
                score=None, attempt_id=None,
            )
            db.add(ap)
            db.flush()  # get ap.id
            token_payload = {
                "sub": f"approve:{ap.id}",
                "type": "approval",
                "exp": datetime.now(timezone.utc) + timedelta(days=7),
            }
            ap.approval_token = jwt.encode(token_payload, SECRET_KEY, algorithm=ALGORITHM)

            # Notify the crew member
            notify(
                db, user.id, "passed", "Course completed",
                f"You have completed all lessons in {course.title}. Your certificate is pending admin approval.",
                f"/my-courses",
            )
            auto_completed = True

    db.commit()
    return {"ok": True, "completed": done, "autoCompleted": auto_completed}


class AssessmentSubmission(BaseModel):
    answers: list[int]


@app.post("/api/courses/{course_id}/assessment")
def submit_assessment(course_id: str, sub: AssessmentSubmission,
                      user: models.User = Depends(get_current_user), db: Session = Depends(get_db)):
    course = db.query(models.Course).filter_by(id=course_id).first()
    if not course:
        raise HTTPException(404, "Course not found")
    require_enrollment(db, user, course_id)
    questions = sorted(course.questions, key=lambda q: q.order)
    if len(sub.answers) != len(questions):
        raise HTTPException(400, "Answer count does not match question count")

    lid = user.id
    p = get_progress(db, lid, course_id)
    already_passed = bool(p and p.passed)

    used = (db.query(models.Attempt)
            .filter_by(learner_id=user.id, course_id=course_id).count())
    if course.max_attempts and not already_passed:
        if used >= course.max_attempts:
            raise HTTPException(
                403, "No assessment attempts remaining. Contact your training officer.")

    correct = sum(1 for q, a in zip(questions, sub.answers) if a == q.answer)
    score = round(correct / len(questions) * 100)
    passed = score >= course.pass_mark

    # full audit trail: one row per submission
    attempt = models.Attempt(learner_id=user.id, course_id=course_id,
                          score=score, passed=passed, answers=sub.answers)
    db.add(attempt)
    db.flush()  # to get attempt.id
    
    # Progress holds the best result so a later worse attempt can't un-pass
    if not p:
        p = models.Progress(learner_id=lid, course_id=course_id, completed_chapters=[])
        db.add(p)
    if p.score is None or score > p.score:
        p.score = score
    p.passed = bool(p.passed) or passed

    if passed:
        # Check if they already have a certificate (e.g. they passed before and retook it)
        cert = None
        existing_cert = db.query(models.Certificate).filter_by(learner_id=lid, course_id=course.id).first()
        if existing_cert:
            cert = cert_dict(existing_cert, user, course)
            notify(db, user.id, "passed", "Assessment passed",
                   f"You passed {course.title} with {score}%. Your certificate is already available.",
                   f"/course/{course.slug}/certificate")
        else:
            # Queue for approval
            ap = models.AssessmentApproval(
                learner_id=user.id, course_id=course.id, score=score, attempt_id=attempt.id
            )
            db.add(ap)
            db.flush()
            
            # Generate the approval token immediately
            from .auth import SECRET_KEY, ALGORITHM
            token_payload = {
                "sub": f"approve:{ap.id}",
                "type": "approval",
                "exp": datetime.now(timezone.utc) + timedelta(days=7)
            }
            ap.approval_token = jwt.encode(token_payload, SECRET_KEY, algorithm=ALGORITHM)
            
            notify(db, user.id, "passed", "Assessment passed",
                   f"You passed {course.title} with {score}%. Your result is pending admin approval for certificate generation.",
                   f"/course/{course.slug}/assessment")
    else:
        cert = None
        notify(db, user.id, "failed", "Assessment not passed",
               f"You scored {score}% on {course.title}. You can retry the assessment.",
               f"/course/{course.slug}/assessment")
    db.commit()

    can_retry = False
    if not passed:
        if course.max_attempts is None:
            can_retry = True
        else:
            can_retry = (used + 1) < course.max_attempts

    # always return review, but mask the correct answer if they failed and can retry
    show_answers = passed or not can_retry
    review = [{
        "q": q.prompt,
        "options": q.options,
        "correct": q.answer if (show_answers or a == q.answer) else None,
        "chosen": a,
        "isCorrect": a == q.answer,
        "explain": q.explain if (show_answers or a == q.answer) else None,
    } for q, a in zip(questions, sub.answers)]

    return {
        "score": score, "passed": passed, "correct": correct,
        "total": len(questions), "certificate": cert, "review": review,
        "attemptsUsed": used + 1, "maxAttempts": course.max_attempts,
        "certPending": passed and not cert
    }


def _course_code(slug: str) -> str:
    """Abbreviate a course slug for certificate IDs: initials of each
    hyphen-separated word (e.g. "mental-health" -> "MH",
    "cargo-operations" -> "CO"). Single-word slugs are kept as-is since an
    initial alone would be too short to be meaningful."""
    words = slug.split("-")
    if len(words) < 2:
        return slug.upper()
    return "".join(w[0] for w in words if w).upper()


def issue_certificate(db, user, course):
    lid = user.id
    existing = (db.query(models.Certificate)
                .filter_by(learner_id=lid, course_id=course.id).first())
    if existing:
        return cert_dict(existing, user, course)
    year = datetime.now(timezone.utc).year
    
    seq_record = models.CertificateSequence()
    db.add(seq_record)
    db.flush() # flush to get seq_record.id
    
    seq = seq_record.id
    cid = f"OZ-{_course_code(course.slug)}-{year}-{seq:04d}"
    cert = models.Certificate(id=cid, learner_id=lid, course_id=course.id,
                              score=get_progress(db, lid, course.id).score)
    db.add(cert)
    db.commit()
    return cert_dict(cert, user, course)


def cert_dict(cert, user, course):
    return {
        "id": cert.id, "learner": user.full_name, "rank": user.rank, "ppNo": user.pp_no,
        "course": course.title, "score": cert.score,
        "issued": cert.issued_at.strftime("%d %B %Y") if cert.issued_at else None,
    }


@app.get("/api/courses/{course_id}/certificate")
def get_certificate(course_id: str, user: models.User = Depends(get_current_user), db: Session = Depends(get_db)):
    course = db.query(models.Course).filter_by(id=course_id).first()
    if not course:
        raise HTTPException(404, "Course not found")
    cert = (db.query(models.Certificate)
            .filter_by(learner_id=user.id, course_id=course_id).first())
    if not cert:
        raise HTTPException(404, "No certificate — assessment not passed yet")
    return cert_dict(cert, user, course)


def cert_pdf_data(cert, user, course):
    c = course.cert or {}
    topics = c.get("topics")
    if not topics:
        # Fallback to chapter titles if no specific topics are set
        unwanted = {
            "introduction", "summary", "conclusion", "quiz", "assessment", 
            "final assessment", "why", "why?", "how", "how?", "what", "what?", 
            "overview", "agenda", "objectives"
        }
        topics = [
            ch.title for ch in course.chapters 
            if ch.kind != "quiz" 
            and ch.title.lower().strip() not in unwanted
            and not ch.title.lower().strip().startswith("slide ")
        ]

    return {
        "id": cert.id,
        "learner": user.full_name,
        "ppNo": user.pp_no,
        "titleUpper": c.get("titleUpper") or course.title.upper(),
        "topics": topics,
        "issued": cert.issued_at.strftime("%d %B %Y") if cert.issued_at else "",
        "location": os.getenv("CERT_LOCATION", "Chennai"),
        "photoPath": os.path.join(UPLOAD_DIR, "photos", f"{user.id}.jpg"),
        "verifyUrl": f"{PUBLIC_BASE_URL}/verify/{cert.id}",
    }


from fastapi.responses import HTMLResponse
import jwt
from .auth import SECRET_KEY, ALGORITHM

def _styled_html_response(title: str, message: str, is_success: bool = True):
    color = "#15a34a" if is_success else "#dc2626"
    icon = "&#10003;" if is_success else "&#10005;"
    return HTMLResponse(f"""
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>Ozellar Marine Training</title>
        <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600&display=swap" rel="stylesheet">
        <style>
            body {{
                font-family: 'Inter', sans-serif;
                background-color: #f6f7f9;
                color: #16181d;
                display: flex;
                align-items: center;
                justify-content: center;
                height: 100vh;
                margin: 0;
            }}
            .card {{
                background: #ffffff;
                padding: 40px;
                border-radius: 12px;
                box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1), 0 2px 4px -1px rgba(0, 0, 0, 0.06);
                text-align: center;
                max-width: 440px;
                width: 90%;
            }}
            .icon-circle {{
                width: 64px;
                height: 64px;
                background-color: {color}15;
                color: {color};
                border-radius: 50%;
                display: flex;
                align-items: center;
                justify-content: center;
                font-size: 32px;
                margin: 0 auto 24px;
            }}
            h2 {{
                margin: 0 0 12px;
                font-size: 20px;
                font-weight: 600;
            }}
            p {{
                margin: 0 0 24px;
                color: #5c626d;
                font-size: 15px;
                line-height: 1.5;
            }}
            .btn {{
                display: inline-block;
                background-color: #2f6fed;
                color: white;
                text-decoration: none;
                padding: 10px 20px;
                border-radius: 8px;
                font-weight: 500;
                font-size: 14px;
                transition: background-color 0.2s;
            }}
            .btn:hover {{
                background-color: #215dd6;
            }}
        </style>
    </head>
    <body>
        <div class="card">
            <div class="icon-circle">{icon}</div>
            <h2>{title}</h2>
            <p>{message}</p>
        </div>
    </body>
    </html>
    """)

@app.get("/api/approve", response_class=HTMLResponse)
def approve_assessment(token: str, db: Session = Depends(get_db)):
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        if payload.get("type") != "approval":
            raise ValueError("Invalid token type")
        sub = payload.get("sub", "")
        if not sub.startswith("approve:"):
            raise ValueError()
        ap_id = sub.split(":")[1]
    except Exception:
        return _styled_html_response("Link Expired", "This approval link is invalid or has expired.", False)

    ap = db.get(models.AssessmentApproval, ap_id)
    if not ap or ap.status != "pending":
        return _styled_html_response("Already Processed", "This approval has already been processed or does not exist.", False)

    user = db.get(models.User, ap.learner_id)
    course = db.query(models.Course).filter_by(id=ap.course_id).first()
    
    # Generate certificate
    cert_info = issue_certificate(db, user, course)
    ap.status = "approved"
    ap.decided_at = datetime.now(timezone.utc)
    
    notify(db, user.id, "certificate", "Certificate Ready",
           f"Your certificate for {course.title} has been approved.",
           f"/course/{course.slug}/certificate")
    db.commit()

    # Generate PDF and email crew
    cert_obj = db.get(models.Certificate, cert_info["id"])
    pdf_bytes = build_certificate_pdf(cert_pdf_data(cert_obj, user, course))
    if user.email:
        email_service.send_approval_email(user.email, user.full_name, course.title, pdf_bytes, cert_info["id"])

    return _styled_html_response(
        "Assessment Approved", 
        f"You have successfully approved {user.full_name} for {course.title}. Their certificate has been generated and dispatched.", 
        True
    )


@app.get("/api/preview-certificate")
def preview_assessment_certificate(token: str, db: Session = Depends(get_db)):
    """Generates a preview PDF of the certificate for an admin reviewing an approval digest email."""
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        if payload.get("type") != "approval":
            raise ValueError("Invalid token type")
        sub = payload.get("sub", "")
        if not sub.startswith("approve:"):
            raise ValueError()
        ap_id = sub.split(":")[1]
    except Exception:
        raise HTTPException(400, "Invalid or expired token")

    ap = db.get(models.AssessmentApproval, ap_id)
    if not ap:
        raise HTTPException(404, "Approval request not found")

    user = db.get(models.User, ap.learner_id)
    course = db.query(models.Course).filter_by(id=ap.course_id).first()
    
    class MockCert:
        id = "PREVIEW-ONLY"
        issued_at = datetime.now(timezone.utc)
        
    pdf = build_certificate_pdf(cert_pdf_data(MockCert(), user, course))
    filename = f"PREVIEW_{user.full_name.replace(' ', '_')}_{course.slug}.pdf"
    return Response(content=pdf, media_type="application/pdf",
                    headers={"Content-Disposition": f'inline; filename="{filename}"'})


@app.get("/api/reject", response_class=HTMLResponse)
def reject_assessment(token: str, db: Session = Depends(get_db)):
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        if payload.get("type") != "approval":
            raise ValueError("Invalid token type")
        sub = payload.get("sub", "")
        if not sub.startswith("approve:"):
            raise ValueError()
        ap_id = sub.split(":")[1]
    except Exception:
        return _styled_html_response("Link Expired", "This rejection link is invalid or has expired.", False)

    ap = db.get(models.AssessmentApproval, ap_id)
    if not ap or ap.status != "pending":
        return _styled_html_response("Already Processed", "This approval has already been processed or does not exist.", False)

    user = db.get(models.User, ap.learner_id)
    course = db.query(models.Course).filter_by(id=ap.course_id).first()
    
    ap.status = "rejected"
    ap.decided_at = datetime.now(timezone.utc)
    
    notify(db, user.id, "failed", "Assessment Rejected",
           f"Your passing score for {course.title} was reviewed but not approved. Please retry.",
           f"/course/{course.slug}/assessment")
           
    # Update progress so they can retake
    p = db.query(models.Progress).filter_by(learner_id=user.id, course_id=course.id).first()
    if p:
        p.passed = False
        
    db.commit()

    if user.email:
        email_service.send_rejection_email(user.email, user.full_name, course.title)

    return _styled_html_response(
        "Assessment Rejected", 
        f"You have rejected the assessment for {user.full_name} for {course.title}. They have been notified to retry.", 
        True
    )

@app.post("/api/crew/photo")
async def upload_crew_photo(file: UploadFile = File(...), user: models.User = Depends(get_current_user)):
    from PIL import Image
    if user.role != "learner":
        raise HTTPException(403, "Only crew members can upload photos")
    if not file.content_type or not file.content_type.startswith("image/"):
        raise HTTPException(400, "File must be an image")
    
    content = await file.read()
    
    try:
        image = Image.open(io.BytesIO(content))
        image.verify() # verify it's an image
    except Exception:
        raise HTTPException(400, "Invalid image file")
    
    # Need to reopen because verify() moves the file pointer
    image = Image.open(io.BytesIO(content))
    
    photos_dir = os.path.join(UPLOAD_DIR, "photos")
    os.makedirs(photos_dir, exist_ok=True)
    
    file_path = os.path.join(photos_dir, f"{user.id}.jpg")
    image.convert("RGB").save(file_path, format="JPEG", quality=85)
        
    return {"status": "success"}



@app.get("/api/courses/{course_id}/certificate.pdf")
def get_certificate_pdf(course_id: str, user: models.User = Depends(get_current_user),
                        db: Session = Depends(get_db)):
    course = db.query(models.Course).filter_by(id=course_id).first()
    if not course:
        raise HTTPException(404, "Course not found")
    cert = (db.query(models.Certificate)
            .filter_by(learner_id=user.id, course_id=course_id).first())
    if not cert:
        raise HTTPException(404, "No certificate — assessment not passed yet")
    pdf = build_certificate_pdf(cert_pdf_data(cert, user, course))
    return Response(content=pdf, media_type="application/pdf",
                    headers={"Content-Disposition": f'inline; filename="{cert.id}.pdf"'})


@app.get("/api/admin/users/{user_id}/courses/{course_id}/certificate.pdf")
def admin_get_crew_certificate_pdf(
    user_id: str, course_id: str,
    request: Request,
    token: Optional[str] = None,
    dl: Optional[str] = None,
    db: Session = Depends(get_db)
):
    """Admin-only: download/view any crew member's issued certificate PDF.
    Accepts JWT via Authorization header OR ?token= query param (for direct <a href> links).
    """
    from .auth import SECRET_KEY, ALGORITHM
    # Resolve token: prefer header, fall back to query param
    auth_header = request.headers.get("Authorization", "")
    raw_token = auth_header.removeprefix("Bearer ").strip() or token
    if not raw_token:
        raise HTTPException(401, "Not authenticated")
    try:
        payload = jwt.decode(raw_token, SECRET_KEY, algorithms=[ALGORITHM])
    except jwt.ExpiredSignatureError:
        raise HTTPException(401, "Session expired")
    except jwt.PyJWTError:
        raise HTTPException(401, "Invalid token")
    if payload.get("type", "session") != "session" or payload.get("role") != "admin":
        raise HTTPException(403, "Admin access required")

    course = db.query(models.Course).filter_by(id=course_id).first()
    if not course:
        raise HTTPException(404, "Course not found")
    learner = db.get(models.User, user_id)
    if not learner:
        raise HTTPException(404, "User not found")
    cert = (db.query(models.Certificate)
            .filter_by(learner_id=user_id, course_id=course_id).first())
    if not cert:
        raise HTTPException(404, "No certificate issued for this crew member and course")
    pdf = build_certificate_pdf(cert_pdf_data(cert, learner, course))
    filename = f"{learner.full_name.replace(' ', '_')}_{course.slug}_{cert.id}.pdf"
    disposition = "attachment" if dl else "inline"
    return Response(content=pdf, media_type="application/pdf",
                    headers={"Content-Disposition": f'{disposition}; filename="{filename}"'})



@app.get("/api/certificates")
def list_certificates(user: models.User = Depends(get_current_user), db: Session = Depends(get_db)):
    certs = db.query(models.Certificate).filter_by(learner_id=user.id).all()
    out = []
    for cert in certs:
        course = db.query(models.Course).filter_by(id=cert.course_id).first()
        if not course:
            continue
        out.append({
            "id": cert.id, "courseId": course.id, "slug": course.slug,
            "course": course.title, "score": cert.score,
            "issued": cert.issued_at.strftime("%d %B %Y") if cert.issued_at else None,
            "pending": False
        })
        
    pending_approvals = db.query(models.AssessmentApproval).filter_by(learner_id=user.id, status="pending").all()
    for ap in pending_approvals:
        course = db.query(models.Course).filter_by(id=ap.course_id).first()
        if not course:
            continue
        out.append({
            "id": f"pending-{ap.id}", "courseId": course.id, "slug": course.slug,
            "course": course.title, "score": ap.score,
            "issued": "Pending",
            "pending": True
        })
    return out


@app.get("/api/verify/{cert_id}")
def verify_certificate(cert_id: str, db: Session = Depends(get_db)):
    """Public certificate verification — no authentication required."""
    cert = db.get(models.Certificate, cert_id)
    if not cert:
        return {"valid": False}
    course = db.query(models.Course).filter_by(id=cert.course_id).first()
    user = db.get(models.User, cert.learner_id)
    return {
        "valid": True, "id": cert.id,
        "holder": user.full_name if user else None,
        "course": course.title if course else None,
        "score": cert.score,
        "issued": cert.issued_at.strftime("%d %B %Y") if cert.issued_at else None,
    }


@app.get("/api/verify/{cert_id}/pdf")
def verify_certificate_pdf(cert_id: str, db: Session = Depends(get_db)):
    """Public — streams the actual certificate PDF for a QR-code scan or
    manual verification link, no authentication required."""
    cert = db.get(models.Certificate, cert_id)
    if not cert:
        raise HTTPException(404, "Certificate not found")
    course = db.query(models.Course).filter_by(id=cert.course_id).first()
    user = db.get(models.User, cert.learner_id)
    if not course or not user:
        raise HTTPException(404, "Certificate not found")
    pdf = build_certificate_pdf(cert_pdf_data(cert, user, course))
    return Response(content=pdf, media_type="application/pdf",
                    headers={"Content-Disposition": f'inline; filename="{cert.id}.pdf"'})


# ======================= NOTIFICATIONS =======================
@app.get("/api/notifications")
def list_notifications(user: models.User = Depends(get_current_user), db: Session = Depends(get_db)):
    rows = (db.query(models.Notification).filter_by(user_id=user.id)
            .order_by(models.Notification.created_at.desc(), models.Notification.id.desc())
            .limit(50).all())
    unread = db.query(models.Notification).filter_by(user_id=user.id, is_read=False).count()
    items = [{
        "id": n.id, "kind": n.kind, "title": n.title, "body": n.body, "link": n.link,
        "isRead": bool(n.is_read),
        "createdAt": n.created_at.isoformat() if n.created_at else None,
    } for n in rows]
    return {"unread": unread, "items": items}


@app.post("/api/notifications/{notif_id}/read")
def mark_notification_read(notif_id: str, user: models.User = Depends(get_current_user),
                           db: Session = Depends(get_db)):
    n = db.get(models.Notification, notif_id)
    if not n or n.user_id != user.id:
        raise HTTPException(404, "Notification not found")
    n.is_read = True
    db.commit()
    return {"ok": True}


@app.post("/api/notifications/read-all")
def mark_all_notifications_read(user: models.User = Depends(get_current_user),
                               db: Session = Depends(get_db)):
    (db.query(models.Notification).filter_by(user_id=user.id, is_read=False)
     .update({models.Notification.is_read: True}))
    db.commit()
    return {"ok": True}


# ======================= ADMIN =======================

@app.get("/api/admin/notifications")
def admin_notifications(admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    """Return pending certificate approvals as notification items for the admin bell."""
    pending = (db.query(models.AssessmentApproval)
               .filter_by(status="pending")
               .order_by(models.AssessmentApproval.created_at.desc())
               .all())
    items = []
    for ap in pending:
        user = db.get(models.User, ap.learner_id)
        course = db.query(models.Course).filter_by(id=ap.course_id).first()
        if not user or not course:
            continue
        items.append({
            "id": ap.id,
            "learnerId": ap.learner_id,
            "courseId": ap.course_id,
            "learnerName": user.full_name,
            "courseName": course.title,
            "createdAt": ap.created_at.isoformat() if ap.created_at else None,
        })
    return {"unread": len(items), "items": items}


class CreateUserRequest(BaseModel):
    role: str                       # 'learner' | 'admin'
    fullName: str
    crewId: str | None = None       # learner
    dob: str | None = None          # learner — 8 digits DDMMYYYY
    rank: str | None = None
    ppNo: str | None = None
    email: str | None = None        # admin
    password: str | None = None     # admin


class UpdateUserRequest(BaseModel):
    isActive: bool | None = None
    fullName: str | None = None
    rank: str | None = None
    ppNo: str | None = None


class AssignRequest(BaseModel):
    courseId: str


def admin_user_view(db, u):
    v = {
        "id": u.id, "role": u.role, "name": u.full_name, "rank": u.rank,
        "crewId": u.crew_id, "email": u.email, "ppNo": u.pp_no,
        "dob": u.date_of_birth.strftime("%d%m%Y") if u.date_of_birth else None,
        "isActive": bool(u.is_active),
    }
    if u.role == "learner":
        v["assignedCount"] = db.query(models.Enrollment).filter_by(learner_id=u.id).count()
        v["passedCount"] = (db.query(models.Progress)
                            .filter_by(learner_id=u.id, passed=True).count())
    return v


@app.get("/api/admin/users")
def admin_list_users(admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    users = db.query(models.User).order_by(models.User.role, models.User.full_name).all()
    return [admin_user_view(db, u) for u in users]


@app.post("/api/admin/users")
def admin_create_user(req: CreateUserRequest, admin: models.User = Depends(require_admin),
                      db: Session = Depends(get_db)):
    name = (req.fullName or "").strip()
    if not name:
        raise HTTPException(400, "Full name is required")

    if req.role == "learner":
        crew_id = (req.crewId or "").strip()
        if not crew_id:
            raise HTTPException(400, "Crew ID is required")
        dob = parse_ddmmyyyy((req.dob or "").strip())
        if dob is None:
            raise HTTPException(400, "Date of birth must be 8 digits (DDMMYYYY)")
        if db.query(models.User).filter_by(crew_id=crew_id).first():
            raise HTTPException(400, "That Crew ID is already in use")
        user = models.User(role="learner", crew_id=crew_id, full_name=name,
                           rank=(req.rank or "").strip() or None, date_of_birth=dob,
                           pp_no=(req.ppNo or "").strip() or None)

    elif req.role == "admin":
        email = (req.email or "").strip().lower()
        if not email:
            raise HTTPException(400, "Email is required")
        if not req.password or len(req.password) < 8:
            raise HTTPException(400, "Password must be at least 8 characters")
        if db.query(models.User).filter_by(email=email).first():
            raise HTTPException(400, "That email is already in use")
        user = models.User(role="admin", email=email, full_name=name,
                           rank=(req.rank or "").strip() or None,
                           password_hash=hash_password(req.password))
    else:
        raise HTTPException(400, "Role must be 'learner' or 'admin'")

    db.add(user)
    db.commit()
    db.refresh(user)
    return admin_user_view(db, user)


@app.patch("/api/admin/users/{user_id}")
def admin_update_user(user_id: str, req: UpdateUserRequest,
                      admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    user = db.get(models.User, user_id)
    if not user:
        raise HTTPException(404, "User not found")
    if req.isActive is not None:
        if user.id == admin.id and req.isActive is False:
            raise HTTPException(400, "You cannot deactivate your own account")
        user.is_active = req.isActive
    if req.fullName is not None:
        user.full_name = req.fullName.strip()
    if req.rank is not None:
        user.rank = req.rank.strip() or None
    if req.ppNo is not None:
        user.pp_no = req.ppNo.strip() or None
    db.commit()
    return admin_user_view(db, user)


@app.post("/api/admin/users/{user_id}/courses/{course_id}/approve")
def admin_inline_approve(user_id: str, course_id: str,
                         admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    """Inline approval of a pending certificate from the Admin Report page."""
    ap = db.query(models.AssessmentApproval).filter_by(
        learner_id=user_id, course_id=course_id, status="pending"
    ).first()
    if not ap:
        raise HTTPException(404, "No pending approval found for this learner/course")
    user = db.get(models.User, user_id)
    course = db.query(models.Course).filter_by(id=course_id).first()
    cert_info = issue_certificate(db, user, course)
    ap.status = "approved"
    ap.decided_at = datetime.now(timezone.utc)
    notify(db, user.id, "certificate", "Certificate Ready",
           f"Your certificate for {course.title} has been approved.",
           f"/course/{course.slug}/certificate")
    db.commit()
    cert_obj = db.get(models.Certificate, cert_info["id"])
    pdf_bytes = build_certificate_pdf(cert_pdf_data(cert_obj, user, course))
    if user.email:
        email_service.send_approval_email(user.email, user.full_name, course.title, pdf_bytes, cert_info["id"])
    return {"ok": True, "certId": cert_info["id"]}


@app.post("/api/admin/users/{user_id}/enrollments")

def admin_assign(user_id: str, req: AssignRequest,
                 admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    user = db.get(models.User, user_id)
    if not user or user.role != "learner":
        raise HTTPException(404, "Learner not found")
    if not db.query(models.Course).filter_by(id=req.courseId).first():
        raise HTTPException(404, "Course not found")
    exists = (db.query(models.Enrollment)
              .filter_by(learner_id=user_id, course_id=req.courseId).first())
    if not exists:
        course = db.query(models.Course).filter_by(id=req.courseId).first()
        db.add(models.Enrollment(learner_id=user_id, course_id=req.courseId,
                                 assigned_by=admin.id))
        notify(db, user_id, "assigned", "New course assigned",
               f"{course.title} has been assigned to you.",
               f"/course/{course.slug}")
        db.commit()
    return {"ok": True}


@app.delete("/api/admin/users/{user_id}/enrollments/{course_id}")
def admin_unassign(user_id: str, course_id: str,
                   admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    # removing an assignment hides the course from the learner but keeps their
    # progress/attempts/certificate rows for audit
    db.query(models.Enrollment).filter_by(learner_id=user_id, course_id=course_id).delete()
    db.commit()
    return {"ok": True}


# ----- compliance reporting -----
def _report(db):
    courses = db.query(models.Course).order_by(models.Course.order).all()
    learners = (db.query(models.User).filter_by(role="learner")
                .order_by(models.User.full_name).all())
    
    enrollments = db.query(models.Enrollment).all()
    progress_all = db.query(models.Progress).all()
    certs = db.query(models.Certificate).all()
    approvals = db.query(models.AssessmentApproval).filter_by(status="pending").all()
    attempts_counts = db.query(models.Attempt.learner_id, models.Attempt.course_id, func.count(models.Attempt.id)).group_by(models.Attempt.learner_id, models.Attempt.course_id).all()
    
    e_map = {(e.learner_id, e.course_id): e for e in enrollments}
    p_map = {(p.learner_id, p.course_id): p for p in progress_all}
    c_map = {(c.learner_id, c.course_id): c for c in certs}
    ap_map = {(a.learner_id, a.course_id): a for a in approvals}
    a_map = {(l, c): count for l, c, count in attempts_counts}
    
    # Pre-compute total chapter count per course (avoids N+1 queries)
    course_chapter_counts = {c.id: len(c.chapters) for c in courses}
    
    rows = []
    for lr in learners:
        assigned_courses = [c for c in courses if (lr.id, c.id) in e_map]
        cells = {}
        for c in assigned_courses:
            enr = e_map.get((lr.id, c.id))
            prog = p_map.get((lr.id, c.id))
            cert = c_map.get((lr.id, c.id))
            ap = ap_map.get((lr.id, c.id))
            attempts = a_map.get((lr.id, c.id), 0)
            
            if prog and prog.passed:
                status = "passed"
            elif prog and ((prog.completed_chapters and len(prog.completed_chapters) > 0) or prog.score is not None):
                status = "in-progress"
            else:
                status = "assigned"
            
            total_chs = course_chapter_counts.get(c.id, 0)
            done_chs  = len(prog.completed_chapters or []) if prog else 0
            pct       = round(done_chs / total_chs * 100) if total_chs else 0
            
            if cert and cert.issued_at:
                passed_on = cert.issued_at.strftime("%Y-%m-%d")
            elif ap and ap.created_at:
                passed_on = ap.created_at.strftime("%Y-%m-%d")
            else:
                passed_on = None
            
            cells[c.id] = {
                "status": status,
                "score": prog.score if prog else None,
                "startedOn": enr.assigned_at.strftime("%Y-%m-%d") if enr and enr.assigned_at else None,
                "passedOn": passed_on,
                "pendingApproval": bool(ap),
                "attempts": attempts,
                "completionPct":     pct,
                "completedChapters": done_chs,
                "totalChapters":     total_chs,
            }
        
        rows.append({
            "learnerId": lr.id, "name": lr.full_name, "crewId": lr.crew_id,
            "rank": lr.rank, "isActive": bool(lr.is_active), "cells": cells,
        })
    return courses, rows


@app.get("/api/admin/report")
def admin_report(admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    courses, rows = _report(db)
    return {"courses": [{"id": c.id, "title": c.title} for c in courses], "rows": rows}


@app.get("/api/admin/dashboard-stats")
def admin_dashboard_stats(admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    """Rich analytics for the admin dashboard charts."""
    from collections import defaultdict
    from datetime import date, timedelta

    # --- crew breakdown ---
    learners = db.query(models.User).filter_by(role="learner").all()
    rank_counts = defaultdict(int)
    status_counts = defaultdict(int)
    vessel_counts = defaultdict(int)
    for lr in learners:
        rank_counts[lr.rank or "Unspecified"] += 1
        status_counts[lr.emp_status or "Unknown"] += 1
        vessel_counts[lr.current_vessel or "Unassigned"] += 1

    rank_data = sorted(
        [{"rank": k, "count": v} for k, v in rank_counts.items()],
        key=lambda x: -x["count"]
    )[:12]  # top 12 ranks

    status_data = [{"status": k, "count": v} for k, v in status_counts.items() if v > 0]

    vessel_data = sorted(
        [{"vessel": k, "count": v} for k, v in vessel_counts.items() if k != "Unassigned"],
        key=lambda x: -x["count"]
    )[:8]

    # --- course stats ---
    courses = db.query(models.Course).order_by(models.Course.order).all()
    enrollments = db.query(models.Enrollment).all()
    progress_all = db.query(models.Progress).all()
    
    e_map = defaultdict(set)
    for e in enrollments:
        e_map[e.course_id].add(e.learner_id)
        
    p_map = {(p.learner_id, p.course_id): p for p in progress_all}
    
    course_stats = []
    for c in courses:
        enrolled_ids = e_map[c.id]
        total = len(enrolled_ids)
        if total == 0:
            course_stats.append({
                "course": c.title[:30], "courseId": c.id,
                "enrolled": 0, "passed": 0, "inProgress": 0,
                "assigned": 0, "passRate": 0,
            })
            continue
        passed = inprogress = 0
        for lid in enrolled_ids:
            prog = p_map.get((lid, c.id))
            if prog and prog.passed:
                passed += 1
            elif prog and ((prog.completed_chapters and len(prog.completed_chapters) > 0) or prog.score is not None):
                inprogress += 1
        assigned = total - passed - inprogress
        course_stats.append({
            "course": c.title[:30], "courseId": c.id,
            "enrolled": total, "passed": passed,
            "inProgress": inprogress, "assigned": assigned,
            "passRate": round((passed / total) * 100) if total else 0,
        })

    # --- enrollment trend (last 6 months) ---
    today = date.today()
    months = []
    for i in range(5, -1, -1):
        # compute start of that month
        month_date = today.replace(day=1)
        for _ in range(i):
            month_date = (month_date - timedelta(days=1)).replace(day=1)
        months.append(month_date)

    enrollment_trend = []
    for m in months:
        if m.month == 12:
            next_m = m.replace(year=m.year + 1, month=1)
        else:
            next_m = m.replace(month=m.month + 1)
        count = (db.query(models.Enrollment)
                 .filter(models.Enrollment.assigned_at >= datetime(m.year, m.month, 1, tzinfo=timezone.utc),
                         models.Enrollment.assigned_at < datetime(next_m.year, next_m.month, 1, tzinfo=timezone.utc))
                 .count())
        enrollment_trend.append({"month": m.strftime("%b %Y"), "enrollments": count})

    # --- recent certificates ---
    recent_certs = (db.query(models.Certificate)
                    .order_by(models.Certificate.issued_at.desc())
                    .limit(8).all())
    recent_cert_list = []
    for cert in recent_certs:
        user_obj = db.get(models.User, cert.learner_id)
        course_obj = db.query(models.Course).filter_by(id=cert.course_id).first()
        recent_cert_list.append({
            "id": cert.id,
            "learner": user_obj.full_name if user_obj else "Unknown",
            "rank": user_obj.rank if user_obj else None,
            "course": course_obj.title if course_obj else cert.course_id,
            "score": cert.score,
            "issuedAt": cert.issued_at.isoformat() if cert.issued_at else None,
        })

    # --- top-level KPIs ---
    total_crew = len(learners)
    active_crew = sum(1 for u in learners if u.is_active)
    total_enrollments = db.query(models.Enrollment).count()
    total_certs = db.query(models.Certificate).count()
    total_courses = db.query(models.Course).count()
    total_attempts = db.query(models.Attempt).count()
    pass_attempts = db.query(models.Attempt).filter_by(passed=True).count()

    return {
        "kpis": {
            "totalCrew": total_crew,
            "activeCrew": active_crew,
            "totalEnrollments": total_enrollments,
            "totalCertificates": total_certs,
            "totalCourses": total_courses,
            "totalAttempts": total_attempts,
            "passAttempts": pass_attempts,
            "overallPassRate": round((pass_attempts / total_attempts) * 100) if total_attempts else 0,
        },
        "crewByRank": rank_data,
        "crewByStatus": status_data,
        "crewByVessel": vessel_data,
        "courseStats": course_stats,
        "enrollmentTrend": enrollment_trend,
        "recentCertificates": recent_cert_list,
    }


@app.get("/api/admin/report.csv")
def admin_report_csv(admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    courses, rows = _report(db)
    buf = io.StringIO()
    w = csv.writer(buf)
    w.writerow(["Crew ID", "Name", "Rank", "Active", "Course", "Status", "Score", "Completed On"])
    for r in rows:
        active = "yes" if r["isActive"] else "no"
        if not r["cells"]:
            w.writerow([r["crewId"], r["name"], r["rank"] or "", active,
                        "(no courses assigned)", "", "", ""])
            continue
        for c in courses:
            cell = r["cells"].get(c.id)
            if not cell:
                continue
            w.writerow([r["crewId"], r["name"], r["rank"] or "", active, c.title,
                        cell["status"], "" if cell["score"] is None else cell["score"],
                        cell["passedOn"] or ""])
    return Response(content=buf.getvalue(), media_type="text/csv",
                    headers={"Content-Disposition": "attachment; filename=ozellar-compliance-report.csv"})


@app.get("/api/admin/report.xlsx")
def admin_report_xlsx(
    admin: models.User = Depends(require_admin),
    db: Session = Depends(get_db),
    crew_search: str = None,   # filter by crew name / crew_id / rank (case-insensitive substring)
    course_id:   str = None,   # filter to a single course column
    status:      str = None,   # filter rows to those with matching status ('passed'|'in-progress'|'assigned')
):
    """Download the compliance report as a styled Excel workbook.

    Supports the same filters exposed in the UI filter bar so that the
    downloaded Excel always matches exactly what the admin sees on screen.
    """
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    courses, rows = _report(db)

    # ---- apply UI filters (mirrors the frontend filter logic) ----
    if crew_search:
        q = crew_search.strip().lower()
        rows = [r for r in rows if
                q in (r["name"] or "").lower() or
                q in (r["crewId"] or "").lower() or
                q in (r["rank"] or "").lower()]

    if course_id:
        rows = [r for r in rows if course_id in r["cells"]]
        courses = [c for c in courses if c.id == course_id]

    if status:
        if course_id:
            rows = [r for r in rows if r["cells"].get(course_id, {}).get("status") == status]
        else:
            rows = [r for r in rows
                    if any(cell["status"] == status for cell in r["cells"].values())]

    # filter indicator for the sheet subtitle
    filter_note = []
    if crew_search: filter_note.append(f'Crew: "{crew_search}"')
    if course_id:
        c_obj = next((c for c in courses if c.id == course_id), None)
        if c_obj: filter_note.append(f"Course: {c_obj.title}")
    if status:
        filter_note.append(f"Status: {status}")

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Compliance Report"

    # ---- colour palette ----
    HDR_FILL   = PatternFill("solid", fgColor="1E3A5F")   # dark navy
    PASS_FILL  = PatternFill("solid", fgColor="D6F5E3")   # soft green
    WIP_FILL   = PatternFill("solid", fgColor="FFF3CD")   # amber
    ASGN_FILL  = PatternFill("solid", fgColor="F2F4F7")   # light grey
    WHITE_FILL = PatternFill("solid", fgColor="FFFFFF")   # white (for merged crew cells)
    HDR_FONT   = Font(bold=True, color="FFFFFF", size=11)
    BODY_FONT  = Font(size=10)
    BOLD_FONT  = Font(bold=True, size=10)
    thin = Side(style="thin", color="D0D5DD")
    thick_bottom = Side(style="medium", color="475467")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    border_thick_bottom = Border(left=thin, right=thin, top=thin, bottom=thick_bottom)
    center = Alignment(horizontal="center", vertical="center")
    left   = Alignment(horizontal="left",   vertical="center", wrap_text=False)

    # ---- title row (row 1): report heading + filter summary ----
    title_text = "Ozellar Marine — Compliance Report"
    if filter_note:
        title_text += f"  |  Filters: {', '.join(filter_note)}"
    ws.merge_cells("A1:H1")
    tc = ws.cell(row=1, column=1, value=title_text)
    tc.font = Font(bold=True, size=12, color="1E3A5F")
    tc.alignment = Alignment(horizontal="left", vertical="center")
    ws.row_dimensions[1].height = 22

    # ---- sub-title row (row 2): generated timestamp + row count ----
    ws.merge_cells("A2:H2")
    sc = ws.cell(row=2, column=1,
                 value=f"Generated: {datetime.now(timezone.utc).strftime('%d %b %Y %H:%M UTC')}  ·  "
                       f"{len(rows)} crew member(s) shown")
    sc.font = Font(size=9, color="5C626D")
    sc.alignment = Alignment(horizontal="left", vertical="center")
    ws.row_dimensions[2].height = 16

    # ---- column headers (row 3) ----
    headers = ["Crew ID", "Name", "Rank", "Active",
               "Course", "Status", "Score (%)", "Attempts", "Started On", "Completed On"]
    ws.append(headers)
    for ci, h in enumerate(headers, 1):
        cell = ws.cell(row=3, column=ci)
        cell.fill = HDR_FILL
        cell.font = HDR_FONT
        cell.alignment = center
        cell.border = border
    ws.row_dimensions[3].height = 22

    status_labels = {"passed": "Completed", "in-progress": "In Progress", "assigned": "Not Started"}
    row_idx = 4
    for r in rows:
        active = "Yes" if r["isActive"] else "No"
        start_row = row_idx
        
        if not r["cells"]:
            data = [r["crewId"], r["name"], r["rank"] or "", active,
                    "(no courses assigned)", "", "", "", "", ""]
            ws.append(data)
            for ci in range(1, 11):
                c = ws.cell(row=row_idx, column=ci)
                c.fill = ASGN_FILL if ci > 4 else WHITE_FILL
                c.font = BODY_FONT
                c.alignment = left; c.border = border
            row_idx += 1
            continue

        for course in courses:
            cell_data = r["cells"].get(course.id)
            if not cell_data:
                continue
            status = cell_data["status"]
            label  = status_labels.get(status, status)
            score  = cell_data["score"] if cell_data["score"] is not None else ""
            started_on = cell_data.get("startedOn") or ""
            passed_on = cell_data.get("passedOn") or ""
            attempts = cell_data.get("attempts", 0)
            if attempts == 0: attempts = ""

            fill = PASS_FILL if status == "passed" else (WIP_FILL if status == "in-progress" else ASGN_FILL)
            data = [r["crewId"], r["name"], r["rank"] or "", active,
                    course.title, label, score, attempts, started_on, passed_on]
            ws.append(data)
            for ci, val in enumerate(data, 1):
                c = ws.cell(row=row_idx, column=ci)
                # columns 1-4 are the crew info (merged later), keep them white
                c.fill = fill if ci > 4 else WHITE_FILL
                c.font = BOLD_FONT if ci == 2 else BODY_FONT
                c.alignment = left
                c.border = border
            row_idx += 1
            
        # merge cells for crew info if they span multiple rows
        if row_idx - 1 > start_row:
            for ci in range(1, 5):
                ws.merge_cells(start_row=start_row, start_column=ci, end_row=row_idx-1, end_column=ci)
                
        # apply thick bottom border to the last row of this user's block to separate users
        for ci in range(1, 11):
            ws.cell(row=row_idx-1, column=ci).border = border_thick_bottom

    # ---- auto column widths ----
    col_widths = [14, 26, 18, 8, 36, 14, 11, 10, 14, 14]
    for ci, w in enumerate(col_widths, 1):
        ws.column_dimensions[get_column_letter(ci)].width = w

    ws.freeze_panes = "A4"  # freeze title + subtitle + header rows

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return Response(
        content=buf.read(),
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=ozellar-compliance-report.xlsx"},
    )


@app.get("/api/crew/my-report.xlsx")
def crew_my_report_xlsx(status: Optional[str] = None, user: models.User = Depends(get_current_user), db: Session = Depends(get_db)):
    """Crew personal training record as a styled Excel workbook.

    One sheet titled "My Training Record" with columns:
    Course | Status | Chapters Completed | Total Chapters | Progress (%) |
    Score (%) | Grade | Time Taken (days) | Certificate ID | Completed On

    Rows are colour-coded: green = passed, amber = in-progress, grey = not started.
    """
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    if user.role == "admin":
        raise HTTPException(403, "Use /api/admin/report.xlsx for admin reports")

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "My Training Record"

    # ---- colour palette ----
    HDR_FILL  = PatternFill("solid", fgColor="1E3A5F")
    PASS_FILL = PatternFill("solid", fgColor="D6F5E3")
    WIP_FILL  = PatternFill("solid", fgColor="FFF3CD")
    ASGN_FILL = PatternFill("solid", fgColor="F2F4F7")
    HDR_FONT  = Font(bold=True, color="FFFFFF", size=11)
    BODY_FONT = Font(size=10)
    BOLD_FONT = Font(bold=True, size=10)
    thin   = Side(style="thin", color="D0D5DD")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    center = Alignment(horizontal="center", vertical="center")
    left   = Alignment(horizontal="left",   vertical="center")

    # ---- title block ----
    ws.merge_cells("A1:L1")
    title_cell = ws.cell(row=1, column=1,
                         value=f"Training Record — {user.full_name}  ·  {user.rank or ''}")
    title_cell.font = Font(bold=True, size=13, color="1E3A5F")
    title_cell.alignment = Alignment(horizontal="left", vertical="center")
    ws.row_dimensions[1].height = 26

    ws.merge_cells("A2:L2")
    sub_cell = ws.cell(row=2, column=1,
                       value=f"Crew ID: {user.crew_id or 'N/A'}   |   Generated: {datetime.now(timezone.utc).strftime('%d %b %Y')}")
    sub_cell.font = Font(size=10, color="5C626D")
    ws.row_dimensions[2].height = 18

    headers = ["Course", "Status", "Attempts", "Chapters Done", "Total Chapters",
               "Progress (%)", "Score (%)", "Grade", "Started On", "Completed On",
               "Time Taken (days)", "Certificate ID"]
    ws.append(headers)   # row 3
    for ci, h in enumerate(headers, 1):
        c = ws.cell(row=3, column=ci)
        c.fill = HDR_FILL; c.font = HDR_FONT
        c.alignment = center; c.border = border
    ws.row_dimensions[3].height = 22

    assigned_ids = enrolled_course_ids(db, user.id)
    courses = (db.query(models.Course).filter(models.Course.id.in_(assigned_ids))
               .order_by(models.Course.order).all()) if assigned_ids else []

    row_idx = 4
    total_passed = 0
    for c in courses:
        prog = get_progress(db, user.id, c.id)
        cert = (db.query(models.Certificate)
                .filter_by(learner_id=user.id, course_id=c.id).first())

        done_count  = len(prog.completed_chapters or []) if prog else 0
        total_ch    = len(c.chapters)
        pct         = round(done_count / total_ch * 100) if total_ch else 0
        score       = prog.score if prog else None
        passed_flag = bool(prog.passed) if prog else False

        if passed_flag:
            if not cert:
                calc_status = "Pending Approval"
                fill = WIP_FILL
            else:
                calc_status = "Completed"
                fill = PASS_FILL
                total_passed += 1
        elif prog and (done_count > 0 or score is not None):
            calc_status = "In Progress"
            fill = WIP_FILL
        else:
            calc_status = "Not Started"
            fill = ASGN_FILL

        # Apply status filter
        if status == 'completed' and calc_status != "Completed": continue
        if status == 'in-progress' and calc_status not in ["In Progress", "Pending Approval"]: continue
        if status == 'not-started' and calc_status != "Not Started": continue

        # grade letter
        grade = ""
        if score is not None:
            if score >= 90:   grade = "A+"
            elif score >= 80: grade = "A"
            elif score >= 70: grade = "B"
            elif score >= 60: grade = "C"
            else:             grade = "F"

        # time taken: enrollment assigned_at → cert issued_at (days)
        time_days = ""
        started_on = ""
        enr = (db.query(models.Enrollment)
               .filter_by(learner_id=user.id, course_id=c.id).first())
        if enr and enr.assigned_at:
            started_on = enr.assigned_at.strftime("%d %b %Y")
            if cert and cert.issued_at:
                delta = cert.issued_at - enr.assigned_at
                time_days = max(0, delta.days)

        attempts = (db.query(models.Attempt)
                    .filter_by(learner_id=user.id, course_id=c.id).count())
        if attempts == 0: attempts = ""

        cert_id   = cert.id if cert else ""
        if passed_flag and not cert:
            passed_on = "Pending Approval"
        else:
            passed_on = cert.issued_at.strftime("%d %b %Y") if cert and cert.issued_at else ""
        score_str = score if score is not None else ""

        row_data = [c.title, calc_status, attempts, done_count, total_ch,
                    pct, score_str, grade, started_on, passed_on, time_days, cert_id]
        ws.append(row_data)

        for ci, val in enumerate(row_data, 1):
            cell = ws.cell(row=row_idx, column=ci)
            cell.fill = fill
            cell.font = BOLD_FONT if ci == 1 else BODY_FONT
            cell.alignment = left
            cell.border = border
        row_idx += 1

    # ---- summary row ----
    ws.append([])
    row_idx += 1
    summary_row = row_idx
    ws.cell(row=summary_row, column=1, value="SUMMARY").font = Font(bold=True, size=10, color="1E3A5F")
    ws.cell(row=summary_row, column=2,
            value=f"{total_passed} of {len(courses)} courses passed").font = BODY_FONT

    # ---- freeze panes & auto widths ----
    col_widths = [36, 15, 10, 15, 15, 14, 11, 8, 14, 14, 18, 22]
    for ci, w in enumerate(col_widths, 1):
        ws.column_dimensions[get_column_letter(ci)].width = w

    ws.freeze_panes = "A4"  # freeze title + header rows

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return Response(
        content=buf.read(),
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename=ozellar-training-{user.crew_id or user.id}.xlsx"},
    )


# Lets an admin create a brand-new course from scratch: upload a .pptx
# (one slide -> one lesson chapter), upload standalone/attached videos,
# insert non-blocking checkpoint quizzes anywhere in the chapter sequence,
# freely reorder everything, and author the mandatory graded final
# assessment. Scoped to new courses only — the 3 seeded courses stay
# managed via courses_seed.json + reseed.

def _slugify(title: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", title.strip().lower()).strip("-") or "course"


def _unique_course_slug(db, base: str) -> str:
    slug = base
    i = 2
    while db.query(models.Course).filter_by(slug=slug).first():
        slug = f"{base}-{i}"
        i += 1
    return slug


def _next_chapter_order(course) -> int:
    return (max((ch.order for ch in course.chapters), default=-1)) + 1


def _next_chapter_n(course) -> int:
    return (max((ch.n or 0 for ch in course.chapters), default=0)) + 1


def admin_course_summary(course):
    return {
        "id": course.id, "slug": course.slug, "title": course.title,
        "subtitle": course.subtitle, "passMark": course.pass_mark,
        "maxAttempts": course.max_attempts,
        "chapterCount": len(course.chapters), "questionCount": len(course.questions),
    }


def admin_chapter_detail(ch):
    return {
        "id": ch.id, "kind": ch.kind, "n": ch.n, "title": ch.title,
        "image": ch.image, "videos": ch.videos, "order": ch.order,
        "quizQuestions": [{
            "id": qq.id, "q": qq.prompt, "options": qq.options,
            "answer": qq.answer, "explain": qq.explain,
        } for qq in ch.quiz_questions] if ch.kind == "quiz" else [],
    }


class CreateCourseRequest(BaseModel):
    title: str
    subtitle: str | None = None
    icon: str | None = None
    gradient: str | None = None
    durationLabel: str | None = None
    passMark: int = 80
    maxAttempts: int | None = None
    targetRanks: list[str] = []
    targetUsers: list[str] = []


class UpdateCourseRequest(BaseModel):
    title: str
    subtitle: str | None = None
    durationLabel: str | None = None
    passMark: int
    maxAttempts: int | None = None
    targetRanks: list[str] = []
    targetUsers: list[str] = []


class CreateQuizChapterRequest(BaseModel):
    title: str
    afterChapterId: str | None = None


class QuizQuestionIn(BaseModel):
    q: str
    options: list[str]
    answer: int
    explain: str | None = None


class SaveQuizQuestionsRequest(BaseModel):
    questions: list[QuizQuestionIn]


class SaveAssessmentRequest(BaseModel):
    passMark: int
    maxAttempts: int | None = None
    questions: list[QuizQuestionIn]


class ReorderRequest(BaseModel):
    order: list[str]


@app.get("/api/admin/courses")
def admin_list_courses(admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    courses = db.query(models.Course).order_by(models.Course.order).all()
    return [admin_course_summary(c) for c in courses]


@app.post("/api/admin/courses")
def admin_create_course(req: CreateCourseRequest, admin: models.User = Depends(require_admin),
                        db: Session = Depends(get_db)):
    title = req.title.strip()
    if not title:
        raise HTTPException(400, "Title is required")
    c_slug = _unique_course_slug(db, _slugify(title))
    existing_orders = [c.order for c in db.query(models.Course).all()]
    course = models.Course(
        slug=c_slug, title=title, subtitle=req.subtitle,
        icon=req.icon, gradient=req.gradient, duration_label=req.durationLabel,
        status="not-started", pass_mark=req.passMark, max_attempts=req.maxAttempts,
        target_ranks=req.targetRanks, target_users=req.targetUsers,
        order=(max(existing_orders) + 1) if existing_orders else 0,
    )
    db.add(course)
    db.flush()  # Generate the UUID for course.id

    # Auto-enroll matching users immediately
    target_ranks = req.targetRanks or []
    target_users = req.targetUsers or []
    if target_ranks or target_users:
        target_ranks_upper = [r.upper() for r in target_ranks]
        users_to_enroll = db.query(models.User).filter(
            models.User.role == 'learner',
            (func.upper(models.User.rank).in_(target_ranks_upper)) | (models.User.id.in_(target_users))
        ).all()
        for u in users_to_enroll:
            enroll = models.Enrollment(learner_id=u.id, course_id=course.id, assigned_by=admin.id)
            db.add(enroll)

    db.commit()
    db.refresh(course)
    return admin_course_summary(course)


@app.put("/api/admin/courses/{course_id}")
def admin_update_course(course_id: str, req: UpdateCourseRequest, admin: models.User = Depends(require_admin),
                        db: Session = Depends(get_db)):
    course = db.get(models.Course, course_id)
    if not course:
        raise HTTPException(404, "Course not found")
        
    title = req.title.strip()
    if not title:
        raise HTTPException(400, "Title is required")
        
    course.title = title
    course.subtitle = req.subtitle
    course.duration_label = req.durationLabel
    course.pass_mark = req.passMark
    course.max_attempts = req.maxAttempts
    course.target_ranks = req.targetRanks
    
    # Filter out legacy integer IDs (e.g., 294) to prevent database DataError
    valid_users = [str(u) for u in req.targetUsers if isinstance(u, str) and len(str(u)) > 20]
    course.target_users = valid_users
    
    db.commit()
    
    # Auto-enroll matching users immediately who are not already enrolled
    target_ranks = req.targetRanks or []
    target_users = req.targetUsers or []
    if target_ranks or target_users:
        target_ranks_upper = [r.upper() for r in target_ranks]
        users_to_enroll = db.query(models.User).filter(
            models.User.role == 'learner',
            (func.upper(models.User.rank).in_(target_ranks_upper)) | (models.User.id.in_(target_users))
        ).all()
        
        existing_enrollments = db.query(models.Enrollment).filter(models.Enrollment.course_id == course_id).all()
        existing_user_ids = {e.learner_id for e in existing_enrollments}
        
        for u in users_to_enroll:
            if u.id not in existing_user_ids:
                enroll = models.Enrollment(learner_id=u.id, course_id=course_id, assigned_by=admin.id)
                db.add(enroll)
        db.commit()
        
    db.refresh(course)
    return admin_course_summary(course)


@app.get("/api/admin/courses/{course_id}")
def admin_get_course_builder(course_id: str, admin: models.User = Depends(require_admin),
                             db: Session = Depends(get_db)):
    course = db.get(models.Course, course_id)
    if not course:
        raise HTTPException(404, "Course not found")
    return {
        "id": course.id, "slug": course.slug, "title": course.title,
        "subtitle": course.subtitle, "durationLabel": course.duration_label,
        "passMark": course.pass_mark, "maxAttempts": course.max_attempts,
        "targetRanks": course.target_ranks or [],
        "targetUsers": course.target_users or [],
        "chapters": [admin_chapter_detail(ch) for ch in
                     sorted(course.chapters, key=lambda c: c.order)],
        "assessment": {
            "passMark": course.pass_mark, "maxAttempts": course.max_attempts,
            "questions": [{
                "id": q.id, "q": q.prompt, "options": q.options,
                "answer": q.answer, "explain": q.explain,
            } for q in sorted(course.questions, key=lambda q: q.order)],
        },
    }


from fastapi import BackgroundTasks
import time

def process_pptx_background(course_id: str, pptx_path: str, original_filename: str, course_dir: str):
    db = SessionLocal()
    try:
        course = db.query(models.Course).filter(models.Course.id == course_id).first()
        if not course: return

        import string
        import random
        # Use a local temp dir to avoid Windows tempfile locking issues with soffice
        tmp = os.path.join(UPLOAD_DIR, "tmp_" + "".join(random.choices(string.ascii_lowercase + string.digits, k=10)))
        os.makedirs(tmp, exist_ok=True)
        
        try:
            soffice_paths = [
                "soffice",
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                r"C:\LibreOffice\program\soffice.exe"
            ]
            
            success = False
            last_err = None
            env_tmp = os.path.join(tmp, "soffice_env").replace(os.sep, "/")
            
            for sp in soffice_paths:
                try:
                    subprocess.run(
                        [sp, f"-env:UserInstallation=file:///{env_tmp}", "--headless", "--nologo", "--nofirststartwizard", "--convert-to", "pdf", "--outdir", tmp, pptx_path],
                        check=True, capture_output=True, timeout=120,
                    )
                    success = True
                    break
                except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired) as e:
                    last_err = e

            pdf_filename = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
            pdf_path = os.path.join(tmp, pdf_filename)
            has_pdf = success and os.path.exists(pdf_path)
            
            if not has_pdf:
                print(f"Background PPTX conversion to PDF failed or skipped: {last_err}")

            import fitz
            from pptx import Presentation

            try:
                prs = Presentation(pptx_path)
                slide_titles = []
                slide_full_texts = []
                for slide in prs.slides:
                    text = ""
                    full_text_parts = []
                    for shape in slide.shapes:
                        if getattr(shape, "has_text_frame", False):
                            full_text_parts.append(shape.text_frame.text.strip())
                        elif getattr(shape, "shape_type", None) == 6:
                            def _get_group_text(shp):
                                res = []
                                for s in shp.shapes:
                                    if getattr(s, "has_text_frame", False):
                                        res.append(s.text_frame.text.strip())
                                    elif getattr(s, "shape_type", None) == 6:
                                        res.extend(_get_group_text(s))
                                return res
                            full_text_parts.extend(_get_group_text(shape))
                    slide_full_texts.append("\n".join(full_text_parts))

                    title_shape = slide.shapes.title
                    if title_shape and title_shape.has_text_frame:
                        text = title_shape.text_frame.text.strip()
                    
                    if not text:
                        def _get_texts(shapes):
                            res = []
                            for s in shapes:
                                if getattr(s, "shape_type", None) == 6:
                                    res.extend(_get_texts(s.shapes))
                                elif getattr(s, "has_text_frame", False):
                                    t = s.text_frame.text.strip()
                                    if len(t) > 3:
                                        res.append((getattr(s, "top", 0) or 0, t))
                            return res
                        texts = _get_texts(slide.shapes)
                        if texts:
                            texts.sort(key=lambda x: x[0])
                            text = texts[0][1].split("\n")[0].strip()
                            
                    slide_titles.append(text if text else "")
            except Exception as e:
                print(f"Failed to extract text with python-pptx: {e}")
                prs = None
                slide_titles = []
                slide_full_texts = []
            
            import zipfile
            import xml.etree.ElementTree as ET
            
            slide_videos = {}
            try:
                with zipfile.ZipFile(pptx_path, "r") as z:
                    for name in z.namelist():
                        if name.startswith("ppt/slides/_rels/slide") and name.endswith(".xml.rels"):
                            try:
                                slide_num_str = name.split("slide")[2].split(".")[0]
                                slide_idx = int(slide_num_str) - 1
                                
                                rels_data = z.read(name)
                                root = ET.fromstring(rels_data)
                                ns = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}
                                for rel in root.findall("r:Relationship", ns):
                                    target = rel.get("Target")
                                    if target and target.startswith("../media/") and target.lower().endswith(".mp4"):
                                        media_path = "ppt/" + target[3:]
                                        if media_path in z.namelist():
                                            media_data = z.read(media_path)
                                            basename = os.path.basename(target)
                                            vid_filename = f"slide{slide_num_str}_{basename}"
                                            vid_path = os.path.join(course_dir, vid_filename)
                                            with open(vid_path, "wb") as vf:
                                                vf.write(media_data)
                                            final_vid_path = compress_video(vid_path)
                                            if final_vid_path == vid_path:
                                                try:
                                                    from qtfaststart import processor
                                                    processor.process(vid_path, vid_path + ".tmp")
                                                    os.replace(vid_path + ".tmp", vid_path)
                                                except Exception:
                                                    pass
                                            storage.save(course_id, vid_filename, final_vid_path)
                                            vid_url = f"/api/uploads/{course_id}/{vid_filename}"
                                            if slide_idx not in slide_videos:
                                                slide_videos[slide_idx] = set()
                                            slide_videos[slide_idx].add(vid_url)
                            except Exception as e:
                                print(f"Error parsing {name}: {e}")
            except Exception as e:
                print(f"Error parsing PPTX zip for videos: {e}")

            doc = fitz.open(pdf_path) if has_pdf else None
            num_slides = len(doc) if doc else (len(prs.slides) if prs else 0)
            
            if num_slides == 0:
                print("No slides found in PPTX or PDF")
                return

            start_n = max((ch.n or 0 for ch in course.chapters), default=0)
            next_order = _next_chapter_order(course)
            
            for i in range(num_slides):
                n = start_n + i + 1
                
                img_url = None
                if doc:
                    page = doc[i]
                    pix = page.get_pixmap(dpi=150)
                    filename = f"slide{n}.png"
                    img_path = os.path.join(course_dir, filename)
                    pix.save(img_path)
                    storage.save(course_id, filename, img_path)
                    img_url = f"/api/uploads/{course_id}/{filename}"
                
                slide_title = slide_titles[i] if i < len(slide_titles) and slide_titles[i] else f"Slide {n}"
                vids = list(slide_videos.get(i, set()))
                
                full_text = slide_full_texts[i] if i < len(slide_full_texts) else ""
                is_quiz = ("quiz" in slide_title.lower()) or ("A. " in full_text and "B. " in full_text)
                ch_kind = "quiz" if is_quiz else "lesson"
                
                quiz_questions = []
                if is_quiz:
                    import re
                    lines = full_text.split("\n")
                    prompt_lines = []
                    options = []
                    ans_idx = 0
                    for line in lines:
                        line_s = line.strip()
                        if not line_s: continue
                        
                        ans_match = re.search(r"answer:\s*([A-E])", line_s, re.IGNORECASE)
                        if ans_match:
                            ans_idx = ord(ans_match.group(1).upper()) - ord("A")
                            continue

                        if re.match(r"^[A-E]\.", line_s):
                            options.append(line_s)
                        else:
                            if not options and line_s.lower() != "quiz":
                                prompt_lines.append(line_s)
                    
                    if options:
                        quiz_questions = [models.ChapterQuestion(
                            prompt=" ".join(prompt_lines).strip() or slide_title,
                            options=options,
                            answer=ans_idx,
                            explain="Auto-extracted from slide."
                        )]
                
                import uuid
                ch = models.Chapter(
                    id=f"{course_id}-slide-{n}-{uuid.uuid4().hex[:8]}", course_id=course_id, n=n,
                    title=slide_title,
                    sections=[], videos=vids, order=next_order + i, kind=ch_kind,
                    image=img_url,
                    quiz_questions=quiz_questions
                )
                db.add(ch)
            
            if doc:
                doc.close()
            db.commit()
            print(f"Background PPTX processing finished for {course_id}")

        finally:
            import shutil
            try:
                shutil.rmtree(tmp)
            except OSError:
                pass
            db.close()
            
    except Exception as e:
        print(f"Background PPTX processing failed entirely: {e}")
        import traceback
        traceback.print_exc()
    finally:
        db.close()
        if os.path.exists(pptx_path):
            try:
                os.remove(pptx_path)
            except OSError:
                pass

@app.post("/api/admin/courses/{course_id}/upload-pptx")
async def upload_course_pptx(course_id: str, background_tasks: BackgroundTasks,
                             file: UploadFile = File(...),
                             admin: models.User = Depends(require_admin),
                             db: Session = Depends(get_db)):
    course = db.query(models.Course).filter(models.Course.id == course_id).first()
    if not course:
        raise HTTPException(404, "Course not found")
    if not (file.filename or "").lower().endswith((".pptx", ".pptm")):
        raise HTTPException(400, "File must be a .pptx or .pptm")

    course_dir = os.path.join(UPLOAD_DIR, course_id)
    os.makedirs(course_dir, exist_ok=True)
    
    pptx_filename = f"pending_{int(time.time())}_{file.filename}"
    pptx_path = os.path.join(course_dir, pptx_filename)
    with open(pptx_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    background_tasks.add_task(process_pptx_background, course_id, pptx_path, file.filename, course_dir)
    return {"message": "Processing started in background."}

@app.post("/api/admin/courses/{course_id}/upload-video")
async def admin_upload_video(course_id: str, file: UploadFile = File(...),
                             chapterId: str | None = Form(None), title: str | None = Form(None),
                             admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    course = db.get(models.Course, course_id)
    if not course:
        raise HTTPException(404, "Course not found")

    course_dir = os.path.join(UPLOAD_DIR, course_id)
    os.makedirs(course_dir, exist_ok=True)
    ext = os.path.splitext(file.filename or "")[1] or ".mp4"
    prefix = f"/api/uploads/{course_id}/video"
    existing = sum(1 for ch in course.chapters for v in (ch.videos or []) if v.startswith(prefix))
    filename = f"video{existing + 1}{ext}"

    raw_path = os.path.join(course_dir, f"_upload_{uuid.uuid4().hex[:8]}{ext}")
    with open(raw_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    if ext.lower() == ".mp4":
        final_path = compress_video(raw_path)
        if final_path == raw_path:
            # ffmpeg unavailable/failed — fall back to just relocating the moov
            # atom so playback can still start immediately.
            try:
                from qtfaststart import processor
                processor.process(raw_path, raw_path + ".tmp")
                os.replace(raw_path + ".tmp", raw_path)
            except Exception as e:
                print("qtfaststart failed:", e)
    else:
        final_path = raw_path

    storage.save(course_id, filename, final_path)
    url = f"/api/uploads/{course_id}/{filename}"

    if chapterId:
        ch = db.get(models.Chapter, chapterId)
        if not ch or ch.course_id != course_id:
            raise HTTPException(404, "Chapter not found")
        ch.videos = [*(ch.videos or []), url]
        db.commit()
        return admin_chapter_detail(ch)

    ch = models.Chapter(
        id=f"{course_id}-video-{_next_chapter_n(course)}-{uuid.uuid4().hex[:8]}", course_id=course_id,
        n=_next_chapter_n(course), title=(title or "Video").strip() or "Video",
        sections=[], videos=[url], order=_next_chapter_order(course), kind="lesson",
    )
    db.add(ch)
    db.commit()
    db.refresh(ch)
    return admin_chapter_detail(ch)


@app.post("/api/admin/courses/{course_id}/quiz-chapters")
def admin_create_quiz_chapter(course_id: str, req: CreateQuizChapterRequest,
                              admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    course = db.get(models.Course, course_id)
    if not course:
        raise HTTPException(404, "Course not found")
    chapters = sorted(course.chapters, key=lambda c: c.order)

    if req.afterChapterId:
        idx = next((i for i, ch in enumerate(chapters) if ch.id == req.afterChapterId), None)
        if idx is None:
            raise HTTPException(404, "Reference chapter not found")
        insert_at = idx + 1
    else:
        insert_at = len(chapters)

    quiz = models.Chapter(
        id=f"{course_id}-quiz-{_next_chapter_n(course)}-{uuid.uuid4().hex[:8]}", course_id=course_id,
        n=_next_chapter_n(course), title=req.title.strip() or "Quiz",
        sections=[], videos=[], kind="quiz", order=0,
    )
    chapters.insert(insert_at, quiz)
    db.add(quiz)
    for i, ch in enumerate(chapters):
        ch.order = i
        ch.n = i + 1
    db.commit()
    db.refresh(quiz)
    return admin_chapter_detail(quiz)


@app.put("/api/admin/courses/{course_id}/chapters/{chapter_id}/quiz-questions")
def admin_save_quiz_questions(course_id: str, chapter_id: str, req: SaveQuizQuestionsRequest,
                              admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    ch = db.get(models.Chapter, chapter_id)
    if not ch or ch.course_id != course_id:
        raise HTTPException(404, "Chapter not found")
    if ch.kind != "quiz":
        raise HTTPException(400, "This chapter is not a quiz")
    db.query(models.ChapterQuestion).filter_by(chapter_id=chapter_id).delete()
    for i, q in enumerate(req.questions):
        db.add(models.ChapterQuestion(
            chapter_id=chapter_id, prompt=q.q, options=q.options,
            answer=q.answer, explain=q.explain, order=i,
        ))
    db.commit()
    db.refresh(ch)
    return admin_chapter_detail(ch)


@app.put("/api/admin/courses/{course_id}/reorder")
def admin_reorder_chapters(course_id: str, req: ReorderRequest,
                           admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    course = db.get(models.Course, course_id)
    if not course:
        raise HTTPException(404, "Course not found")
    if set(req.order) != {ch.id for ch in course.chapters}:
        raise HTTPException(400, "Order must include exactly the course's current chapters")
    by_id = {ch.id: ch for ch in course.chapters}
    for i, cid in enumerate(req.order):
        by_id[cid].order = i
        by_id[cid].n = i + 1
    db.commit()
    return {"ok": True}


@app.delete("/api/admin/courses/{course_id}/chapters/{chapter_id}")
def admin_delete_chapter(course_id: str, chapter_id: str,
                         admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    ch = db.get(models.Chapter, chapter_id)
    if not ch or ch.course_id != course_id:
        raise HTTPException(404, "Chapter not found")
    prefix = f"/api/uploads/{course_id}/"
    files_to_remove = [p for p in [ch.image, *(ch.videos or [])] if p and p.startswith(prefix)]
    db.query(models.ChapterQuestion).filter_by(chapter_id=chapter_id).delete()
    db.delete(ch)
    db.flush()
    # Re-index remaining chapters so 'n' and 'order' remain contiguous
    course = db.get(models.Course, course_id)
    if course:
        remaining = sorted(course.chapters, key=lambda c: c.order)
        for i, r in enumerate(remaining):
            r.order = i
            r.n = i + 1
    db.commit()
    for p in files_to_remove:
        storage.delete(course_id, os.path.basename(p))
    return {"ok": True}


@app.put("/api/admin/courses/{course_id}/assessment")
def admin_save_assessment(course_id: str, req: SaveAssessmentRequest,
                          admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    course = db.get(models.Course, course_id)
    if not course:
        raise HTTPException(404, "Course not found")
    course.pass_mark = req.passMark
    course.max_attempts = req.maxAttempts
    db.query(models.Question).filter_by(course_id=course_id).delete()
    for i, q in enumerate(req.questions):
        db.add(models.Question(
            course_id=course_id, prompt=q.q, options=q.options,
            answer=q.answer, explain=q.explain, order=i,
        ))
    db.commit()
    return admin_get_course_builder(course_id, admin, db)


# ============================================================
# SCREENING / ENTRANCE TEST — TEST-TAKER ROUTES
# ============================================================

SCREENING_PHOTO_DIR = os.path.join(UPLOAD_DIR, "screening_photos")
os.makedirs(SCREENING_PHOTO_DIR, exist_ok=True)


class ScreeningLoginRequest(BaseModel):
    name: str
    password: str


@app.post("/api/screening/login")
def screening_login(req: ScreeningLoginRequest, db: Session = Depends(get_db)):
    """Test-taker login: full name + password."""
    name = (req.name or "").strip()
    if not name or not req.password:
        raise HTTPException(401, "Name and password are required")
    check_rate_limit(db, f"screening:{normalize_name(name)}")
    candidates = db.query(models.ScreeningCandidate).filter_by(is_active=True).all()
    target = normalize_name(name)
    match = None
    for c in candidates:
        if normalize_name(c.full_name) == target:
            if verify_password(req.password, c.password_hash):
                match = c
                break
    if not match:
        raise HTTPException(401, "Invalid name or password")
    if match.status == "submitted":
        raise HTTPException(403, "You have already submitted this test")
    clear_rate_limit(db, f"screening:{normalize_name(name)}")
    token = create_screening_token(match)
    return {"token": token, "candidate": candidate_public(match)}


@app.get("/api/screening/test")
def screening_get_test(
    candidate: models.ScreeningCandidate = Depends(get_current_candidate),
    db: Session = Depends(get_db),
):
    """Get the full test details for the current candidate (no answer keys)."""
    test = db.get(models.ScreeningTest, candidate.test_id)
    if not test or not test.is_active:
        raise HTTPException(404, "Test not found or inactive")
    attempt = candidate.attempt
    sections_data = []
    for sec in test.sections:
        s = {
            "id": sec.id, "title": sec.title,
            "type": sec.section_type, "order": sec.order,
            "passage": sec.passage,
            "questions": [
                {"id": q.id, "prompt": q.prompt, "options": q.options,
                 "imageUrls": q.image_urls or [], "order": q.order}
                for q in sec.questions
            ],
        }
        sections_data.append(s)
    # Remaining time is computed server-side against the DB clock. `started_at`
    # comes from the DB's now(), so comparing it to the DB's now() again keeps
    # both values on one wall clock — the candidate's browser timezone (and any
    # clock skew) can't shorten or extend the exam.
    remaining_seconds = None
    if attempt and attempt.started_at and not attempt.submitted_at:
        db_now = db.query(func.now()).scalar()
        started = attempt.started_at
        if db_now.tzinfo is not None:
            db_now = db_now.replace(tzinfo=None)
        if started.tzinfo is not None:
            started = started.replace(tzinfo=None)
        elapsed = (db_now - started).total_seconds()
        remaining_seconds = max(0, int(test.timer_minutes * 60 - elapsed))

    return {
        "id": test.id, "title": test.title,
        "timerMinutes": test.timer_minutes,
        "correctScore": test.correct_score,
        "wrongPenalty": test.wrong_penalty,
        "remainingSeconds": remaining_seconds,
        "sections": sections_data,
        "attempt": {
            "startedAt": attempt.started_at.isoformat() if attempt else None,
            "submittedAt": attempt.submitted_at.isoformat() if attempt and attempt.submitted_at else None,
            "status": candidate.status,
            "tabSwitchCount": attempt.tab_switch_count or 0,
            # Last autosaved (or submitted) answers — lets the client resume
            # from the server's copy if local browser storage was lost
            # (device power-off, different browser/device, cleared storage).
            "sectionAnswers": attempt.section_answers or {},
            "personalData": attempt.personal_data or {},
        } if attempt else None,
    }


@app.post("/api/screening/start")
def screening_start(
    candidate: models.ScreeningCandidate = Depends(get_current_candidate),
    db: Session = Depends(get_db),
):
    """Create the attempt row (records start time). Idempotent."""
    if candidate.status == "submitted":
        raise HTTPException(403, "Test already submitted")
    existing = candidate.attempt
    if existing:
        return {"startedAt": existing.started_at.isoformat(), "ok": True}
    attempt = models.ScreeningAttempt(
        candidate_id=candidate.id, test_id=candidate.test_id,
    )
    db.add(attempt)
    candidate.status = "in_progress"
    db.commit()
    db.refresh(attempt)
    return {"startedAt": attempt.started_at.isoformat(), "ok": True}


@app.post("/api/screening/tab-switch")
def screening_tab_switch(
    candidate: models.ScreeningCandidate = Depends(get_current_candidate),
    db: Session = Depends(get_db),
):
    """Record one instance of the candidate leaving the exam tab/window."""
    if candidate.status == "submitted":
        raise HTTPException(403, "Test already submitted")
    attempt = candidate.attempt
    if not attempt:
        raise HTTPException(400, "Test not started")
    attempt.tab_switch_count = (attempt.tab_switch_count or 0) + 1
    db.commit()
    return {"tabSwitchCount": attempt.tab_switch_count}


@app.post("/api/screening/photo")
async def screening_upload_photo(
    file: UploadFile = File(...),
    candidate: models.ScreeningCandidate = Depends(get_current_candidate),
    db: Session = Depends(get_db),
):
    """Upload passport-size photo for the candidate."""
    from PIL import Image
    data = await file.read()
    try:
        img = Image.open(io.BytesIO(data))
        img.verify()
    except Exception:
        raise HTTPException(400, "Invalid image file")
    dest = os.path.join(SCREENING_PHOTO_DIR, f"{candidate.id}.jpg")
    img2 = Image.open(io.BytesIO(data)).convert("RGB")
    img2.thumbnail((600, 600))
    img2.save(dest, "JPEG", quality=85)
    if candidate.attempt:
        candidate.attempt.photo_path = dest
    db.commit()
    return {"ok": True, "hasPhoto": True}


@app.get("/api/screening/photo")
def screening_get_photo(
    candidate: models.ScreeningCandidate = Depends(get_current_candidate),
):
    """Serve the signed-in candidate's own identity photo (exam header/sidebar)."""
    path = os.path.join(SCREENING_PHOTO_DIR, f"{candidate.id}.jpg")
    if not os.path.exists(path):
        raise HTTPException(404, "No photo on record")
    with open(path, "rb") as fh:
        return Response(fh.read(), media_type="image/jpeg",
                        headers={"Cache-Control": "private, max-age=300"})


class ScreeningSubmitRequest(BaseModel):
    personal_data: Optional[dict] = None
    section_answers: dict  # {section_id: [int | null]}


@app.post("/api/screening/submit")
def screening_submit(
    req: ScreeningSubmitRequest,
    candidate: models.ScreeningCandidate = Depends(get_current_candidate),
    db: Session = Depends(get_db),
):
    """Grade and finalise the test. Returns the scored result."""
    if candidate.status == "submitted":
        raise HTTPException(403, "Test already submitted")
    test = db.get(models.ScreeningTest, candidate.test_id)
    if not test:
        raise HTTPException(404, "Test not found")

    attempt = candidate.attempt
    if not attempt:
        attempt = models.ScreeningAttempt(candidate_id=candidate.id, test_id=candidate.test_id)
        db.add(attempt)
        db.flush()

    correct = wrong = unanswered = 0
    section_results = []
    for sec in test.sections:
        if sec.section_type == "personal_data":
            continue
        answers = req.section_answers.get(sec.id, [])
        sq = 0; swrong = 0; sunanswered = 0
        for i, q in enumerate(sec.questions):
            chosen = answers[i] if i < len(answers) else None
            if chosen is None:
                sunanswered += 1
            elif chosen == q.answer:
                sq += 1
            else:
                swrong += 1
        correct += sq
        wrong += swrong
        unanswered += sunanswered
        section_results.append({
            "sectionId": sec.id, "title": sec.title,
            "correct": sq, "wrong": swrong, "unanswered": sunanswered,
            "total": len(sec.questions),
        })

    raw_score = correct * test.correct_score - wrong * test.wrong_penalty

    attempt.personal_data = req.personal_data
    attempt.section_answers = req.section_answers
    attempt.score = raw_score
    attempt.correct_count = correct
    attempt.wrong_count = wrong
    attempt.unanswered_count = unanswered
    attempt.submitted_at = datetime.now(timezone.utc)

    candidate.status = "submitted"
    if req.personal_data:
        mob = req.personal_data.get("mobile") or req.personal_data.get("mobileNumber")
        if mob and not candidate.mobile_number:
            candidate.mobile_number = mob
    db.commit()
    return {
        "score": raw_score, "correct": correct, "wrong": wrong,
        "unanswered": unanswered,
        "correctScore": test.correct_score,
        "wrongPenalty": test.wrong_penalty,
        "sectionResults": section_results,
        "submittedAt": attempt.submitted_at.isoformat(),
    }


@app.post("/api/screening/autosave")
def screening_autosave(
    req: ScreeningSubmitRequest,
    candidate: models.ScreeningCandidate = Depends(get_current_candidate),
    db: Session = Depends(get_db),
):
    """Periodically persist in-progress answers so a device power-off,
    crash, or connection drop mid-test doesn't lose the candidate's work.
    Never grades or finalises anything — that only happens in /submit."""
    if candidate.status == "submitted":
        raise HTTPException(403, "Test already submitted")
    attempt = candidate.attempt
    if not attempt:
        raise HTTPException(400, "Test not started")
    attempt.section_answers = req.section_answers
    if req.personal_data:
        attempt.personal_data = req.personal_data
    db.commit()
    return {"ok": True, "savedAt": datetime.now(timezone.utc).isoformat()}


@app.get("/api/screening/result")
def screening_result(
    candidate: models.ScreeningCandidate = Depends(get_current_candidate),
    db: Session = Depends(get_db),
):
    """Get the result for the current candidate (after submission)."""
    if candidate.status != "submitted":
        raise HTTPException(400, "Test not yet submitted")
    attempt = candidate.attempt
    if not attempt:
        raise HTTPException(404, "Attempt not found")
    test = db.get(models.ScreeningTest, candidate.test_id)

    # Attemptable total comes from the graded sections, so the result page can
    # show a real "x / max" and percentage rather than dividing by nothing.
    total_questions = sum(
        len(s.questions) for s in (test.sections or [])
        if s.section_type != "personal_data"
    ) if test else 0
    max_score = total_questions * (test.correct_score if test else 4)

    time_taken_minutes = None
    if attempt.started_at and attempt.submitted_at:
        started, ended = attempt.started_at, attempt.submitted_at
        if (started.tzinfo is None) != (ended.tzinfo is None):
            started = started.replace(tzinfo=None)
            ended = ended.replace(tzinfo=None)
        time_taken_minutes = max(0, round((ended - started).total_seconds() / 60))

    return {
        "score": attempt.score, "correct": attempt.correct_count,
        "wrong": attempt.wrong_count, "unanswered": attempt.unanswered_count,
        "correctScore": test.correct_score if test else 4,
        "wrongPenalty": test.wrong_penalty if test else 1,
        "totalQuestions": total_questions,
        "maxScore": max_score,
        "timeTakenMinutes": time_taken_minutes,
        "fullName": candidate.full_name,
        "testTitle": test.title if test else None,
        "submittedAt": attempt.submitted_at.isoformat() if attempt.submitted_at else None,
        "startedAt": attempt.started_at.isoformat() if attempt.started_at else None,
    }


# ============================================================
# ADMIN — SCREENING MANAGEMENT
# ============================================================

class CreateTestRequest(BaseModel):
    title: str
    timerMinutes: int = 80
    correctScore: int = 4
    wrongPenalty: int = 1


class SectionRequest(BaseModel):
    id: Optional[str] = None
    title: str
    section_type: str = "mcq"
    passage: Optional[str] = None
    order: int = 0


class QuestionRequest(BaseModel):
    id: Optional[str] = None
    prompt: str
    options: list[str]
    answer: int
    imageUrls: Optional[list[str]] = None
    order: int = 0


class BulkQuestionsRequest(BaseModel):
    questions: list[QuestionRequest]


class CreateCandidateRequest(BaseModel):
    fullName: str
    password: str
    testId: str
    mobileNumber: Optional[str] = None


def serialize_test(t, include_sections=False, db=None):
    candidate_count = len(t.candidates) if t.candidates else 0
    submitted_count = sum(1 for c in (t.candidates or []) if c.status == "submitted")
    total_questions = sum(
        len(s.questions) for s in (t.sections or []) if s.section_type == "mcq"
    )
    data = {
        "id": t.id, "title": t.title, "timerMinutes": t.timer_minutes,
        "correctScore": t.correct_score, "wrongPenalty": t.wrong_penalty,
        "isActive": t.is_active, "createdAt": t.created_at.isoformat() if t.created_at else None,
        "candidateCount": candidate_count, "submittedCount": submitted_count,
        "totalQuestions": total_questions,
    }
    if include_sections:
        data["sections"] = [
            {
                "id": s.id, "title": s.title, "type": s.section_type,
                "passage": s.passage, "order": s.order,
                "questions": [
                    {"id": q.id, "prompt": q.prompt, "options": q.options,
                     "answer": q.answer, "imageUrls": q.image_urls or [], "order": q.order}
                    for q in s.questions
                ],
            }
            for s in (t.sections or [])
        ]
    return data


@app.get("/api/admin/screening/tests")
def admin_list_tests(admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    tests = db.query(models.ScreeningTest).order_by(models.ScreeningTest.created_at.desc()).all()
    return [serialize_test(t) for t in tests]


@app.post("/api/admin/screening/tests")
def admin_create_test(req: CreateTestRequest,
                       admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    test = models.ScreeningTest(
        title=req.title, timer_minutes=req.timerMinutes,
        correct_score=req.correctScore, wrong_penalty=req.wrongPenalty,
    )
    db.add(test)
    db.flush()
    # Auto-create the default personal data section
    db.add(models.ScreeningSection(
        test_id=test.id, title="Personal Details",
        section_type="personal_data", order=0,
    ))
    db.commit()
    db.refresh(test)
    return serialize_test(test, include_sections=True)


@app.get("/api/admin/screening/tests/{test_id}")
def admin_get_test(test_id: str,
                   admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    test = db.get(models.ScreeningTest, test_id)
    if not test:
        raise HTTPException(404, "Test not found")
    return serialize_test(test, include_sections=True)


@app.patch("/api/admin/screening/tests/{test_id}")
def admin_update_test(test_id: str, req: CreateTestRequest,
                       admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    test = db.get(models.ScreeningTest, test_id)
    if not test:
        raise HTTPException(404, "Test not found")
    test.title = req.title
    test.timer_minutes = req.timerMinutes
    test.correct_score = req.correctScore
    test.wrong_penalty = req.wrongPenalty
    db.commit()
    return serialize_test(test, include_sections=True)


@app.patch("/api/admin/screening/tests/{test_id}/toggle")
def admin_toggle_test(test_id: str,
                       admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    test = db.get(models.ScreeningTest, test_id)
    if not test:
        raise HTTPException(404, "Test not found")
    test.is_active = not test.is_active
    db.commit()
    return {"isActive": test.is_active}


@app.delete("/api/admin/screening/tests/{test_id}")
def admin_delete_test(test_id: str,
                       admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    test = db.get(models.ScreeningTest, test_id)
    if not test:
        raise HTTPException(404, "Test not found")
    db.delete(test)
    db.commit()
    return {"ok": True}


@app.post("/api/admin/screening/tests/{test_id}/sections")
def admin_add_section(test_id: str, req: SectionRequest,
                       admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    test = db.get(models.ScreeningTest, test_id)
    if not test:
        raise HTTPException(404, "Test not found")
    # Determine next order
    existing_orders = [s.order for s in test.sections]
    next_order = max(existing_orders) + 1 if existing_orders else 1
    sec = models.ScreeningSection(
        test_id=test_id, title=req.title,
        section_type=req.section_type, passage=req.passage,
        order=next_order,
    )
    db.add(sec)
    db.commit()
    db.refresh(test)
    return serialize_test(test, include_sections=True)


@app.patch("/api/admin/screening/tests/{test_id}/sections/{section_id}")
def admin_update_section(test_id: str, section_id: str, req: SectionRequest,
                          admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    sec = db.get(models.ScreeningSection, section_id)
    if not sec or sec.test_id != test_id:
        raise HTTPException(404, "Section not found")
    sec.title = req.title
    sec.passage = req.passage
    db.commit()
    test = db.get(models.ScreeningTest, test_id)
    return serialize_test(test, include_sections=True)


@app.delete("/api/admin/screening/tests/{test_id}/sections/{section_id}")
def admin_delete_section(test_id: str, section_id: str,
                          admin: models.User = Depends(require_admin), db: Session = Depends(get_db)):
    sec = db.get(models.ScreeningSection, section_id)
    if not sec or sec.test_id != test_id:
        raise HTTPException(404, "Section not found")
    if sec.section_type == "personal_data":
        raise HTTPException(400, "Cannot delete the personal data section")
    db.delete(sec)
    db.commit()
    test = db.get(models.ScreeningTest, test_id)
    return serialize_test(test, include_sections=True)


@app.put("/api/admin/screening/tests/{test_id}/sections/{section_id}/questions")
def admin_save_section_questions(
    test_id: str, section_id: str, req: BulkQuestionsRequest,
    admin: models.User = Depends(require_admin), db: Session = Depends(get_db),
):
    """Replace all questions in a section."""
    sec = db.get(models.ScreeningSection, section_id)
    if not sec or sec.test_id != test_id:
        raise HTTPException(404, "Section not found")
    db.query(models.ScreeningQuestion).filter_by(section_id=section_id).delete()
    for i, q in enumerate(req.questions):
        db.add(models.ScreeningQuestion(
            section_id=section_id, prompt=q.prompt, options=q.options,
            answer=q.answer, image_urls=q.imageUrls, order=i,
        ))
    db.commit()
    test = db.get(models.ScreeningTest, test_id)
    return serialize_test(test, include_sections=True)


@app.post("/api/admin/screening/upload-question-image")
async def admin_upload_screening_question_image(
    file: UploadFile = File(...),
    admin: models.User = Depends(require_admin),
):
    """Upload an image for a screening MCQ question (figure/diagram-based
    reasoning questions). Returns the URL to save into the question's
    imageUrl — goes through the same storage backend as everything else
    (local disk in dev, Azure Blob in prod), unlike the always-local
    candidate identity photos."""
    from PIL import Image
    data = await file.read()
    try:
        img = Image.open(io.BytesIO(data))
        img.verify()
    except Exception:
        raise HTTPException(400, "Invalid image file")
    img2 = Image.open(io.BytesIO(data)).convert("RGB")
    img2.thumbnail((1000, 1000))
    filename = f"{uuid.uuid4().hex}.jpg"
    tmp_path = os.path.join(UPLOAD_DIR, f"_tmp_{filename}")
    img2.save(tmp_path, "JPEG", quality=88)
    storage.save("screening-questions", filename, tmp_path)
    return {"imageUrl": f"/api/uploads/screening-questions/{filename}"}


# --- Candidates ---

@app.get("/api/admin/screening/candidates")
def admin_list_candidates(
    test_id: Optional[str] = None,
    admin: models.User = Depends(require_admin), db: Session = Depends(get_db)
):
    q = db.query(models.ScreeningCandidate)
    if test_id:
        q = q.filter_by(test_id=test_id)
    candidates = q.order_by(models.ScreeningCandidate.created_at.desc()).all()
    result = []
    for c in candidates:
        test = db.get(models.ScreeningTest, c.test_id)
        att = c.attempt
        result.append({
            "id": c.id, "fullName": c.full_name, "mobileNumber": c.mobile_number,
            "status": c.status, "isActive": c.is_active,
            "testId": c.test_id, "testTitle": test.title if test else "—",
            "createdAt": c.created_at.isoformat() if c.created_at else None,
            "startedAt": att.started_at.isoformat() if att and att.started_at else None,
            "submittedAt": att.submitted_at.isoformat() if att and att.submitted_at else None,
            "score": att.score if att else None,
            "tabSwitchCount": att.tab_switch_count or 0 if att else 0,
        })
    return result


@app.post("/api/admin/screening/candidates")
def admin_create_candidate(
    req: CreateCandidateRequest,
    admin: models.User = Depends(require_admin), db: Session = Depends(get_db)
):
    test = db.get(models.ScreeningTest, req.testId)
    if not test:
        raise HTTPException(404, "Test not found")
    c = models.ScreeningCandidate(
        test_id=req.testId, full_name=req.fullName.strip(),
        password_hash=hash_password(req.password),
        mobile_number=req.mobileNumber,
    )
    db.add(c)
    db.commit()
    db.refresh(c)
    return {
        "id": c.id, "fullName": c.full_name, "mobileNumber": c.mobile_number,
        "status": c.status, "testId": c.test_id, "testTitle": test.title,
        "createdAt": c.created_at.isoformat() if c.created_at else None,
        "startedAt": None, "submittedAt": None, "score": None,
    }


@app.patch("/api/admin/screening/candidates/{cand_id}")
def admin_update_candidate(
    cand_id: str, req: dict,
    admin: models.User = Depends(require_admin), db: Session = Depends(get_db)
):
    c = db.get(models.ScreeningCandidate, cand_id)
    if not c:
        raise HTTPException(404, "Candidate not found")
    if "isActive" in req:
        c.is_active = req["isActive"]
    if "fullName" in req and req["fullName"].strip():
        c.full_name = req["fullName"].strip()
    if "mobileNumber" in req:
        c.mobile_number = req["mobileNumber"]
    if "password" in req and req["password"]:
        c.password_hash = hash_password(req["password"])
    if "testId" in req and req["testId"] and req["testId"] != c.test_id:
        new_test = db.get(models.ScreeningTest, req["testId"])
        if not new_test:
            raise HTTPException(404, "Test not found")
        # switching tests invalidates any progress/result on the old one
        if c.attempt:
            db.delete(c.attempt)
        c.test_id = req["testId"]
        c.status = "pending"
    db.commit()
    return {"ok": True}


@app.delete("/api/admin/screening/candidates/{cand_id}")
def admin_delete_candidate(
    cand_id: str,
    admin: models.User = Depends(require_admin), db: Session = Depends(get_db)
):
    c = db.get(models.ScreeningCandidate, cand_id)
    if not c:
        raise HTTPException(404, "Candidate not found")
    db.delete(c)
    db.commit()
    return {"ok": True}


# --- Results ---

@app.get("/api/admin/screening/results")
def admin_screening_results(
    test_id: Optional[str] = None,
    admin: models.User = Depends(require_admin), db: Session = Depends(get_db)
):
    q = db.query(models.ScreeningAttempt)
    if test_id:
        q = q.filter_by(test_id=test_id)
    attempts = q.order_by(models.ScreeningAttempt.submitted_at.desc()).all()
    rows = []
    for att in attempts:
        c = db.get(models.ScreeningCandidate, att.candidate_id)
        test = db.get(models.ScreeningTest, att.test_id)
        if not c:
            continue
        # Time taken
        time_taken = None
        if att.started_at and att.submitted_at:
            delta = att.submitted_at - att.started_at
            time_taken = round(delta.total_seconds() / 60, 1)
        # Total possible score
        total_q = sum(
            len(s.questions) for s in (test.sections if test else [])
            if s.section_type == "mcq"
        )
        max_score = total_q * (test.correct_score if test else 4)
        rows.append({
            "candidateId": c.id, "fullName": c.full_name,
            "mobileNumber": c.mobile_number or att.personal_data.get("mobile", "") if att.personal_data else (c.mobile_number or ""),
            "testId": att.test_id, "testTitle": test.title if test else "—",
            "startedAt": att.started_at.isoformat() if att.started_at else None,
            "submittedAt": att.submitted_at.isoformat() if att.submitted_at else None,
            "timeTakenMinutes": time_taken,
            "correct": att.correct_count or 0,
            "wrong": att.wrong_count or 0,
            "unanswered": att.unanswered_count or 0,
            "score": att.score or 0,
            "maxScore": max_score,
            "totalQuestions": total_q,
            "personalData": att.personal_data,
            "tabSwitchCount": att.tab_switch_count or 0,
        })
    return rows


@app.get("/api/admin/screening/results.xlsx")
def admin_screening_results_xlsx(
    test_id: Optional[str] = None,
    admin: models.User = Depends(require_admin), db: Session = Depends(get_db)
):
    """Download screening results as a styled Excel workbook."""
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    rows_data = admin_screening_results(test_id=test_id, admin=admin, db=db)

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Screening Results"

    # Header style
    header_fill = PatternFill("solid", fgColor="1A2744")
    header_font = Font(bold=True, color="FFFFFF", size=11)
    center = Alignment(horizontal="center", vertical="center")
    thin = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin"),
    )

    headers = [
        "Name", "Mobile Number", "Test", "Start Time", "Submit Time",
        "Time Taken (min)", "Correct", "Wrong", "Unanswered", "Score",
        "Remark", "Shortlisted",
    ]
    ws.append(headers)
    for col_idx, _ in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col_idx)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = center
        cell.border = thin
    ws.row_dimensions[1].height = 22

    alt_fill = PatternFill("solid", fgColor="EEF2FF")
    for row_i, r in enumerate(rows_data, 2):
        def fmt_dt(iso):
            if not iso:
                return ""
            try:
                from datetime import datetime as dt
                return dt.fromisoformat(iso.replace("Z", "+00:00")).strftime("%d %b %Y %H:%M")
            except Exception:
                return iso

        ws.append([
            r["fullName"], r["mobileNumber"] or "", r["testTitle"],
            fmt_dt(r["startedAt"]), fmt_dt(r["submittedAt"]),
            r["timeTakenMinutes"] or "",
            r["correct"], r["wrong"], r["unanswered"],
            r["score"],
            "",  # Remark     — to be filled manually by admin
            "",  # Shortlisted — to be filled manually by admin
        ])
        fill = alt_fill if row_i % 2 == 0 else None
        for col_idx in range(1, len(headers) + 1):
            cell = ws.cell(row=row_i, column=col_idx)
            cell.border = thin
            cell.alignment = Alignment(vertical="center")
            if fill:
                cell.fill = fill

    # Column widths
    col_widths = [28, 18, 28, 20, 20, 18, 10, 10, 12, 10, 15, 30]
    for i, w in enumerate(col_widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=screening-results.xlsx"},
    )
